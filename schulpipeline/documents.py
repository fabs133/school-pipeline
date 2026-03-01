"""Document role classification and template filling.

School projects often involve multiple documents with different roles:

  Template:    "Fill this in, don't change the layout"
               → Projektantrag form, PPTX slide deck with structure
               → Output: same file with fields filled in

  Source:      "Extract information from this"
               → Teacher instructions, requirement lists, grading criteria
               → Output: structured data for filling templates / answering tasks

  Constraint:  "Use this as rules for the output"
               → Page limits, formatting requirements, grading rubric
               → Output: validation rules applied to generated content

  Reference:   "Background material, use if relevant"
               → Textbook excerpts, example solutions, law texts
               → Output: cited in answers where applicable

The critical flow for Projektantrag-style tasks:
  1. Classify each input document's role
  2. Parse templates to find fillable fields/placeholders
  3. Extract requirements from source documents
  4. Generate content that FITS the template constraints (length, format)
  5. Fill template in-place, preserving layout

For DOCX templates: python-docx reads and modifies the actual file.
For PPTX templates: python-pptx reads and modifies the actual file.
The layout, styles, fonts, spacing — everything stays. Only text content changes.
"""

from __future__ import annotations

import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from .stages.base import BaseStage
from .stages.intake import _parse_json_response

_logger = logging.getLogger("schulpipeline.documents")

# ============================================================
# Document Role Model
# ============================================================


@dataclass
class ClassifiedDocument:
    """A document with its classified role."""

    id: str  # "doc_01"
    filename: str  # "Projektantrag.docx"
    role: str  # template | source | constraint | reference
    content_type: str  # docx | pptx | pdf | text | image
    content: str  # Extracted text content
    fields: list[TemplateField] | None = None  # For templates: detected fields
    extracted_info: dict[str, Any] = field(default_factory=dict)  # For sources
    constraints: list[str] = field(default_factory=list)  # For constraints
    raw_path: str = ""  # Original file path if available


@dataclass
class TemplateField:
    """A fillable field detected in a template document."""

    id: str  # "field_01"
    label: str  # "Projektbezeichnung"
    field_type: str  # text | paragraph | table_cell | placeholder
    location: str  # Where in the document (for human reference)
    current_value: str = ""  # What's currently there (placeholder text, empty, etc.)
    max_length: int | None = None  # Character limit (estimated from space)
    constraints: list[str] = field(default_factory=list)  # "must fit on one line", etc.


@dataclass
class FilledTemplate:
    """Result of filling a template."""

    template_id: str
    output_path: str
    fields_filled: list[dict[str, str]]  # [{field_id, label, value}]
    fields_skipped: list[dict[str, str]]  # [{field_id, label, reason}]
    warnings: list[str] = field(default_factory=list)


# ============================================================
# Document Classifier Stage
# ============================================================

CLASSIFY_PROMPT = """\
Du bist ein Dokument-Klassifikator für Schulprojekte. Du erhältst eine Liste von
Dokumenten und musst die Rolle jedes Dokuments bestimmen.

Rollen:
- "template": Dokument das ausgefüllt werden soll. Layout DARF NICHT verändert werden.
  Erkennbar an: Leerfelder, Platzhalter, Formularstruktur, "bitte ausfüllen"
- "source": Dokument enthält Informationen/Anforderungen die extrahiert werden sollen.
  Erkennbar an: Aufgabenstellungen, Anforderungslisten, Projektbeschreibungen
- "constraint": Dokument enthält Regeln/Einschränkungen für die Bearbeitung.
  Erkennbar an: Bewertungsschemata, Formatvorgaben, Abgabekriterien
- "reference": Hintergrundmaterial das bei Bedarf genutzt werden kann.
  Erkennbar an: Lehrbuchauszüge, Beispiellösungen, Gesetzestexte

Antworte mit validem JSON:
{
  "documents": [
    {
      "id": "doc_01",
      "filename": "Projektantrag.docx",
      "role": "template",
      "reasoning": "Formular mit Leerfeldern, soll ausgefüllt zurückgegeben werden",
      "fields": [
        {
          "id": "field_01",
          "label": "Projektbezeichnung",
          "field_type": "text",
          "location": "Zeile 3",
          "current_value": "",
          "max_length": 80,
          "constraints": ["einzeilig"]
        }
      ],
      "extracted_info": {}
    },
    {
      "id": "doc_02",
      "filename": "Anforderungen.txt",
      "role": "source",
      "reasoning": "Liste von Anforderungen für das Projekt",
      "fields": null,
      "extracted_info": {
        "requirements": ["muss Datenbank enthalten", "Login-System"],
        "contradictions": ["Lehrer A sagt REST, Lehrer B sagt GraphQL"]
      }
    }
  ],
  "contradictions": [
    {
      "topic": "API-Architektur",
      "source_a": "Anforderungen.txt: REST API",
      "source_b": "Mündliche Aussage: GraphQL",
      "recommendation": "REST verwenden (schriftlich dokumentiert)"
    }
  ]
}

Regeln:
- Jedes Dokument genau EINER Rolle zuordnen
- Bei Templates: alle erkennbaren Felder mit Typ und geschätzter Maximallänge
- Bei Sources: Anforderungen, Widersprüche, und Kerninfos extrahieren
- Widersprüche zwischen Dokumenten explizit markieren
- max_length bei Template-Feldern: Zeichenzahl schätzen basierend auf verfügbarem Platz
- Antworte NUR mit JSON
"""


class ClassifyDocsStage(BaseStage):
    """Classifies input documents by their role in the assignment."""

    name = "classify_docs"
    spec_path = "specs/classify_docs.json"

    async def execute(self, context: dict[str, Any], backend: Any, config: Any) -> dict[str, Any]:
        """Executes the method using the provided context and backend.

        :param context: A dictionary containing necessary information for execution.
        :type context: dict[str, Any]
        :param backend: The backend to be used for processing.
        :type backend: Any
        :param config: Configuration settings for the execution.
        :type config: Any
        :return: A dictionary containing the results of the execution.
        :rtype: dict[str, Any]
        """
        documents = context.get("documents", [])
        _preset = context.get("preset")  # reserved for future use

        # Build document listing for the LLM
        doc_listing = []
        for i, doc in enumerate(documents):
            entry = f"--- Dokument {i + 1}: {doc.get('filename', f'doc_{i}')} ---\n"
            entry += f"Typ: {doc.get('content_type', 'text')}\n"
            entry += f"Inhalt:\n{doc.get('content', '')[:3000]}\n"
            doc_listing.append(entry)

        all_docs = "\n\n".join(doc_listing)

        messages = [
            {"role": "system", "content": CLASSIFY_PROMPT},
            {"role": "user", "content": f"Klassifiziere diese Dokumente:\n\n{all_docs}"},
        ]

        response = await backend.complete(
            stage=self.name,
            messages=messages,
            temperature=0.1,
            max_tokens=8192,
            response_format={"type": "json_object"},
        )

        return _parse_json_response(response.content)


# ============================================================
# Template Filler Stage
# ============================================================

FILL_TEMPLATE_PROMPT = """\
Du bist ein Template-Ausfüller für Schulprojekte. Du erhältst:
1. Ein Template mit leeren Feldern
2. Informationen aus Quell-Dokumenten
3. Einschränkungen (Zeichenlimits, Format)

Deine Aufgabe: Generiere den Text für jedes Feld so, dass er:
- Inhaltlich korrekt und vollständig ist
- In das vorgegebene Zeichenlimit passt
- Dem Sprachstil des Dokuments entspricht (formell/informell)
- Fachlich korrekt ist

{context}

Antworte mit validem JSON:
{{
  "filled_fields": [
    {{
      "field_id": "field_01",
      "label": "Projektbezeichnung",
      "value": "Entwicklung einer webbasierten Lagerverwaltung",
      "fits_constraint": true,
      "char_count": 50
    }}
  ],
  "warnings": [
    "Feld 'Zeitplan' hat nur 200 Zeichen — Detailplan separat empfohlen"
  ]
}}

Regeln:
- JEDES Feld muss einen Wert bekommen (oder eine Begründung warum nicht)
- Zeichenlimit STRIKT einhalten — lieber kürzen als überschreiten
- Bei DIN A4-Constraint: Gesamttext muss auf eine Seite passen
- Deutscher Fachsprachstil, keine Umgangssprache
- Antworte NUR mit JSON
"""


class FillTemplateStage(BaseStage):
    """Fills template fields with generated content that respects constraints."""

    name = "fill_template"
    spec_path = "specs/fill_template.json"

    async def execute(self, context: dict[str, Any], backend: Any, config: Any) -> dict[str, Any]:
        """Executes the method using the provided context, backend, and config.

        :param context: A dictionary containing classification information.
        :type context: dict[str, Any]
        :param backend: The backend to use for execution.
        :type backend: Any
        :param config: Configuration settings for the execution.
        :type config: Any
        :return: A dictionary with the result of the execution or an error message.
        :rtype: dict[str, Any]
        :raises ValueError: If no templates are found in the context.
        """
        classified = context.get("classify_docs", {})
        preset = context.get("preset")

        templates = [d for d in classified.get("documents", []) if d.get("role") == "template"]
        sources = [d for d in classified.get("documents", []) if d.get("role") == "source"]
        constraints = [d for d in classified.get("documents", []) if d.get("role") == "constraint"]

        if not templates:
            return {"error": "Keine Templates gefunden", "filled_templates": []}

        # Consolidate information from sources
        source_info = self._consolidate_sources(sources)
        constraint_rules = self._consolidate_constraints(constraints)
        contradictions = classified.get("contradictions", [])

        results = []
        for template in templates:
            filled = await self._fill_single_template(
                template, source_info, constraint_rules, contradictions, backend, preset
            )
            results.append(filled)

        data = {
            "filled_templates": results,
            "contradictions": contradictions,
            "source_summary": source_info,
        }

        # Write output files (moved from pipeline.py post-processing)
        output_dir = context.get("output_dir")
        if output_dir:
            output_dir = Path(output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)
            for tmpl in results:
                filename = tmpl.get("template_filename", "output")
                fields = tmpl.get("fields_filled", [])
                template_file = context.get("template_files", {}).get(filename)
                if template_file and Path(template_file).exists():
                    out = output_dir / f"filled_{filename}"
                    try:
                        if filename.endswith(".pptx"):
                            apply_to_pptx(template_file, fields, out)
                        else:
                            apply_to_docx(template_file, fields, out)
                        data["file_path"] = str(out)
                        _logger.info(f"Template filled: {out}")
                    except Exception as e:
                        _logger.warning(f"Template fill failed for {filename}: {e}")

        return data

    async def _fill_single_template(
        self, template: dict, source_info: dict, constraints: list[str], contradictions: list, backend: Any, preset: Any
    ) -> dict:
        """Fill a single template document."""
        fields = template.get("fields", []) or []
        if not fields:
            return {
                "template_id": template.get("id", "?"),
                "fields_filled": [],
                "fields_skipped": [],
                "warnings": ["Keine Felder im Template erkannt"],
            }

        # Build context for the LLM
        context_parts = []
        context_parts.append(f"Template: {template.get('filename', '?')}")
        context_parts.append("\nFelder zum Ausfüllen:")
        for f in fields:
            limit = f" (max {f['max_length']} Zeichen)" if f.get("max_length") else ""
            context_parts.append(f"  - {f['label']}: {f['field_type']}{limit}")
            if f.get("constraints"):
                context_parts.append(f"    Einschränkungen: {', '.join(f['constraints'])}")

        context_parts.append("\nVerfügbare Informationen:")
        context_parts.append(json.dumps(source_info, ensure_ascii=False, indent=2)[:3000])

        if constraints:
            context_parts.append("\nFormat-Einschränkungen:")
            for c in constraints:
                context_parts.append(f"  - {c}")

        if contradictions:
            context_parts.append("\nWidersprüche (beachten!):")
            for c in contradictions:
                context_parts.append(f"  - {c.get('topic', '?')}: {c.get('recommendation', 'Unklar')}")

        context_text = "\n".join(context_parts)
        prompt = FILL_TEMPLATE_PROMPT.replace("{context}", context_text)

        if preset:
            prompt += f"\n\nFachkontext: {preset.system_context}"

        messages = [
            {"role": "system", "content": prompt},
            {"role": "user", "content": f"Fülle das Template '{template.get('filename', '')}' aus."},
        ]

        response = await backend.complete(
            stage="fill_template",
            messages=messages,
            temperature=0.2,
            max_tokens=4096,
            response_format={"type": "json_object"},
        )

        try:
            data = _parse_json_response(response.content)
        except ValueError:
            data = {"filled_fields": [], "warnings": ["LLM-Antwort konnte nicht geparst werden"]}

        return {
            "template_id": template.get("id", "?"),
            "template_filename": template.get("filename", "?"),
            "fields_filled": data.get("filled_fields", []),
            "fields_skipped": [],
            "warnings": data.get("warnings", []),
        }

    def _consolidate_sources(self, sources: list[dict]) -> dict:
        """Merge information from all source documents."""
        consolidated = {
            "requirements": [],
            "project_description": "",
            "key_facts": [],
        }
        for src in sources:
            info = src.get("extracted_info", {})
            if info.get("requirements"):
                consolidated["requirements"].extend(info["requirements"])
            if info.get("description"):
                consolidated["project_description"] += info["description"] + "\n"
            for k, v in info.items():
                if k not in ("requirements", "description"):
                    consolidated["key_facts"].append({k: v})
        return consolidated

    def _consolidate_constraints(self, constraint_docs: list[dict]) -> list[str]:
        """Extract constraint rules from constraint documents."""
        rules = []
        for doc in constraint_docs:
            info = doc.get("extracted_info", {})
            if info.get("rules"):
                rules.extend(info["rules"])
            if info.get("format_requirements"):
                rules.extend(info["format_requirements"])
        return rules


# ============================================================
# Template Application — write values back into actual files
# ============================================================


def apply_to_docx(template_path: str | Path, filled_fields: list[dict], output_path: str | Path) -> list[str]:
    """Fill a DOCX template in-place, preserving all formatting.

    Strategy:
      1. Open the template with python-docx
      2. For each field, search for the placeholder/label in the document
      3. Replace the placeholder with the generated value
      4. Save to output_path (never modify the original)

    Returns list of warnings.
    """
    from docx import Document

    template_path = Path(template_path)
    output_path = Path(output_path)

    doc = Document(str(template_path))
    warnings = []

    # Build a lookup: field_label → value
    field_map = {}
    for f in filled_fields:
        label = f.get("label", "")
        value = f.get("value", "")
        if label and value:
            field_map[label] = value

    # Search and replace in paragraphs
    filled_count = 0
    for paragraph in doc.paragraphs:
        for label, value in field_map.items():
            if label in paragraph.text or _normalize(label) in _normalize(paragraph.text):
                _replace_in_paragraph(paragraph, label, value)
                filled_count += 1

    # Search and replace in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for label, value in field_map.items():
                        if label in paragraph.text or _normalize(label) in _normalize(paragraph.text):
                            _replace_in_paragraph(paragraph, label, value)
                            filled_count += 1

    if filled_count < len(field_map):
        warnings.append(
            f"Nur {filled_count}/{len(field_map)} Felder konnten im Dokument gefunden werden. "
            f"Möglicherweise sind die Feldbezeichnungen anders als erwartet."
        )

    doc.save(str(output_path))
    return warnings


def apply_to_pptx(template_path: str | Path, filled_fields: list[dict], output_path: str | Path) -> list[str]:
    """Fill a PPTX template in-place, preserving all formatting.

    Same strategy as DOCX but for PowerPoint slides.
    """
    from pptx import Presentation

    template_path = Path(template_path)
    output_path = Path(output_path)

    prs = Presentation(str(template_path))
    warnings = []

    field_map = {}
    for f in filled_fields:
        label = f.get("label", "")
        value = f.get("value", "")
        if label and value:
            field_map[label] = value

    filled_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for label, value in field_map.items():
                    full_text = "".join(run.text for run in paragraph.runs)
                    if label in full_text or _normalize(label) in _normalize(full_text):
                        _replace_in_pptx_paragraph(paragraph, label, value)
                        filled_count += 1

            # Also check tables in shapes
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for label, value in field_map.items():
                                full_text = "".join(run.text for run in paragraph.runs)
                                if label in full_text or _normalize(label) in _normalize(full_text):
                                    _replace_in_pptx_paragraph(paragraph, label, value)
                                    filled_count += 1

    if filled_count < len(field_map):
        warnings.append(f"Nur {filled_count}/{len(field_map)} Felder im PPTX gefunden.")

    prs.save(str(output_path))
    return warnings


# ============================================================
# Helpers
# ============================================================


def _normalize(text: str) -> str:
    """Normalize text for fuzzy matching: lowercase, strip extra whitespace."""
    return re.sub(r"\s+", " ", text.lower().strip())


def _replace_in_paragraph(paragraph, label: str, value: str) -> None:
    """Replace a label in a DOCX paragraph while preserving formatting.

    Strategy: find which runs contain the label, replace text in those runs.
    If label spans multiple runs (Word splits text unpredictably), join and re-split.
    """
    full_text = paragraph.text

    # Simple case: label is in a single run
    for run in paragraph.runs:
        if label in run.text:
            run.text = run.text.replace(label, value, 1)
            return

    # Complex case: label spans multiple runs
    # Reconstruct text, find position, map back to runs
    if label in full_text:
        # Nuclear option for cross-run labels: clear all runs, set first run's text
        # This preserves the FIRST run's formatting for the whole paragraph
        new_text = full_text.replace(label, value, 1)
        if paragraph.runs:
            # Preserve formatting of first run
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""


def _replace_in_pptx_paragraph(paragraph, label: str, value: str) -> None:
    """Replace a label in a PPTX paragraph while preserving formatting."""
    full_text = "".join(run.text for run in paragraph.runs)

    for run in paragraph.runs:
        if label in run.text:
            run.text = run.text.replace(label, value, 1)
            return

    # Cross-run replacement
    if label in full_text:
        new_text = full_text.replace(label, value, 1)
        if paragraph.runs:
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""


# ============================================================
# Output Preset
# ============================================================

TEMPLATE_FILL_CONFIG = {
    "key": "vorlage",
    "label": "Vorlage ausfüllen",
    "format": "template_fill",
    "stages": ["intake", "classify_docs", "fill_template"],
    "constraints": {
        "template_mode": True,
        "preserve_layout": True,
    },
}
