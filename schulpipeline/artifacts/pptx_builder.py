"""PPTX artifact builder using python-pptx."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..styles import (
    DEFAULT_STYLE,
    DISABLED_VISUAL_SLOTS,
    VisualSlotConfig,
    VisualStyle,
    _lighten,
    _to_rgb,
)

# Visual type labels for placeholders (text-based for reliable rendering)
_TYPE_ICONS = {
    "diagram": "[Diagramm]",
    "photo": "[Foto]",
    "chart": "[Grafik]",
    "icon": "[Symbol]",
    "screenshot": "[Screenshot]",
}

# Headings that indicate a sources/references section
_SOURCES_KEYWORDS = frozenset({
    "quellen", "quellenangaben", "referenzen",
    "literaturverzeichnis", "literatur", "sources", "references",
})


def _blank_layout(prs: Presentation):
    """Find a blank slide layout by name, falling back to index 6 then the last layout."""
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    # Fallback: index 6 (standard in default templates), then last
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[-1]


def _is_sources_section(section: dict, index: int, total: int) -> bool:
    """Detect whether a section is the sources/references slide."""
    heading = section.get("heading", "").lower().strip()
    if index == total - 1:
        return any(kw in heading for kw in _SOURCES_KEYWORDS)
    return heading in _SOURCES_KEYWORDS


def _content_width(prs: Presentation, margin: float) -> int:
    """Compute usable content width in EMU from slide width minus margins."""
    return prs.slide_width - Inches(2 * margin)


def build_pptx(
    synthesis: dict[str, Any],
    output_path: Path,
    visual_style: VisualStyle | None = None,
    visual_config: VisualSlotConfig | None = None,
) -> None:
    """Build a PPTX presentation from synthesis data."""
    if visual_style is None:
        visual_style = DEFAULT_STYLE.visual
    if visual_config is None:
        visual_config = DISABLED_VISUAL_SLOTS

    prs = Presentation()

    # Slide dimensions — exact EMU values to avoid float truncation
    if visual_style.layout.slide_ratio == "4:3":
        prs.slide_width = Emu(9144000)    # 10 inches
    else:
        prs.slide_width = Emu(12192000)   # 13.333 inches (16:9)
    prs.slide_height = Emu(6858000)       # 7.5 inches

    sections = synthesis.get("sections", [])
    title = synthesis.get("title", "Präsentation")
    sources = synthesis.get("sources", [])

    for i, section in enumerate(sections):
        # Skip completely empty sections
        heading = section.get("heading", "").strip()
        content = section.get("content", "").strip()
        bullets = section.get("bullet_points", [])
        if not heading and not content and not bullets:
            continue

        if i == 0:
            _add_title_slide(prs, section, title, visual_style)
        elif _is_sources_section(section, i, len(sections)):
            _add_sources_slide(prs, section, sources, visual_style)
        else:
            _add_content_slide(prs, section, visual_style, visual_config)

    prs.save(str(output_path))


def _add_title_slide(
    prs: Presentation, section: dict, title: str, style: VisualStyle,
) -> None:
    """Dark background title slide."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    margin = style.layout.margin_inches
    cw = _content_width(prs, margin)

    # Dark background (conditional on title_slide_dark)
    bg = slide.background.fill
    bg.solid()
    if style.layout.title_slide_dark:
        bg.fore_color.rgb = _to_rgb(style.colors.bg_dark)
    else:
        bg.fore_color.rgb = _to_rgb(style.colors.bg_light)

    title_color = style.colors.text_light if style.layout.title_slide_dark else style.colors.text_dark
    subtitle_color = style.colors.secondary if style.layout.title_slide_dark else style.colors.text_muted

    # Title
    left, top = Inches(margin), Inches(2.0)
    width, height = cw, Inches(2.0)
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section.get("heading", title)
    p.font.size = Pt(style.fonts.heading_size_pt + 8)
    p.font.bold = True
    p.font.name = style.fonts.heading_family
    p.font.color.rgb = _to_rgb(title_color)
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if section.get("content", "").strip():
        left, top = Inches(margin), Inches(4.2)
        width, height = cw, Inches(1.5)
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section["content"]
        p.font.size = Pt(style.fonts.body_size_pt + 2)
        p.font.name = style.fonts.body_family
        p.font.color.rgb = _to_rgb(subtitle_color)
        p.alignment = PP_ALIGN.LEFT


def _add_content_slide(
    prs: Presentation, section: dict,
    style: VisualStyle, visual_config: VisualSlotConfig,
) -> None:
    """Content slide with heading, bullet points, and optional visual placeholder."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    margin = style.layout.margin_inches
    cw = _content_width(prs, margin)

    # Light background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _to_rgb(style.colors.bg_light)

    # Determine if we have a visual to place
    visuals = section.get("visuals", [])
    has_visual = visual_config.enabled and bool(visuals)
    first_visual = visuals[0] if has_visual else None
    placement = (first_visual or {}).get("placement", visual_config.default_placement)

    # Heading
    left, top = Inches(margin), Inches(0.5)
    width, height = cw, Inches(1.0)
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section.get("heading", "")
    p.font.size = Pt(style.fonts.heading_size_pt)
    p.font.bold = True
    p.font.name = style.fonts.heading_family
    p.font.color.rgb = _to_rgb(style.colors.primary)
    p.alignment = PP_ALIGN.LEFT

    # Bullet points — narrow when visual is on the right
    bullets = section.get("bullet_points", [])
    if bullets:
        left, top = Inches(margin), Inches(1.8)
        if has_visual and placement == "right":
            # Left half: ~50% of slide minus margin
            width = int(prs.slide_width * 0.5) - Inches(margin)
        else:
            width = cw
        height = Inches(5.0)
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        for j, bullet in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"\u2022  {bullet}"
            p.font.size = Pt(style.fonts.bullet_size_pt)
            p.font.name = style.fonts.body_family
            p.font.color.rgb = _to_rgb(style.colors.text_dark)
            p.space_after = Pt(12)
            p.alignment = PP_ALIGN.LEFT

    # Speaker notes
    notes = section.get("speaker_notes")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    # Visual placeholders
    if has_visual:
        for visual in visuals[:visual_config.max_per_slide]:
            _add_visual_placeholder(slide, visual, style, visual_config, prs.slide_width)


def _add_sources_slide(
    prs: Presentation, section: dict, sources: list[str], style: VisualStyle,
) -> None:
    """Sources/references slide."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    margin = style.layout.margin_inches
    cw = _content_width(prs, margin)

    bg = slide.background.fill
    bg.solid()
    if style.layout.title_slide_dark:
        bg.fore_color.rgb = _to_rgb(style.colors.bg_dark)
    else:
        bg.fore_color.rgb = _to_rgb(style.colors.bg_light)

    text_color = style.colors.text_light if style.layout.title_slide_dark else style.colors.text_dark
    source_color = style.colors.secondary if style.layout.title_slide_dark else style.colors.text_muted

    # Heading
    left, top = Inches(margin), Inches(0.5)
    width, height = cw, Inches(1.0)
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = section.get("heading", "Quellen")
    p.font.size = Pt(style.fonts.heading_size_pt)
    p.font.bold = True
    p.font.name = style.fonts.heading_family
    p.font.color.rgb = _to_rgb(text_color)

    # Sources list
    all_sources = sources or section.get("bullet_points", [])
    if all_sources:
        left, top = Inches(margin), Inches(1.8)
        width, height = cw, Inches(5.0)
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        for j, source in enumerate(all_sources):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = source
            p.font.size = Pt(max(style.fonts.bullet_size_pt - 4, 12))
            p.font.name = style.fonts.body_family
            p.font.color.rgb = _to_rgb(source_color)
            p.space_after = Pt(8)


def _add_visual_placeholder(
    slide, visual: dict, style: VisualStyle, config: VisualSlotConfig,
    slide_width: int,
) -> None:
    """Add a styled placeholder box with dashed border for a visual element."""
    placement = visual.get("placement", config.default_placement)
    margin = style.layout.margin_inches

    # Position based on placement — computed from actual slide width
    if placement == "right":
        vis_width = int(slide_width * 0.38)
        left = slide_width - Inches(margin) - vis_width
        top = Inches(1.8)
        height = Inches(4.5)
    else:  # center
        vis_width = _content_width_from(slide_width, margin)
        left = Inches(margin)
        top = Inches(4.0)
        height = Inches(3.0)

    # Rounded rectangle with dashed border, light fill
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, vis_width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _to_rgb(_lighten(style.colors.bg_light))
    shape.line.color.rgb = _to_rgb(style.colors.text_muted)
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.line.width = Pt(1.5)

    # Label + intent + search hint
    icon = _TYPE_ICONS.get(visual.get("type", ""), "[Bild]")
    label = f"{icon}  {visual.get('intent', '')}"
    if config.show_search_hint and visual.get("search_hint"):
        label += f"\n\nSuche: {visual['search_hint']}"

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    p.font.size = Pt(12)
    p.font.name = style.fonts.body_family
    p.font.color.rgb = _to_rgb(style.colors.text_muted)


def _content_width_from(slide_width: int, margin: float) -> int:
    """Compute content width EMU from slide width EMU and margin inches."""
    return slide_width - Inches(2 * margin)
