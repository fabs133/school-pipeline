#!/usr/bin/env python3
"""Standalone test runner — works without pytest.

Run: python3 tests/run_tests.py
"""

import asyncio
import json
import os
import sys
import tempfile
import unittest
from pathlib import Path

# Add project root to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from schulpipeline.backends.base import LLMResponse
from schulpipeline.config import BackendConfig, OutputConfig, PipelineConfig, ResearchConfig, load_config

# --- Mock Backend (inline for standalone use) ---

DEFAULT_RESPONSES = {
    # --- Worksheet-specific (must come BEFORE shorter keys to avoid substring collision) ---
    "Aufgaben-Parser für Schularbeitsblätter": json.dumps(
        {
            "title": "Gewinnverteilung OHG/KG",
            "subject": "Wirtschaft",
            "tasks": [
                {
                    "id": "aufgabe_1a",
                    "label": "Übung 2, Aufgabe 1a",
                    "task_type": "table_fill",
                    "text": "Der Gewinn der Baumeister OHG beträgt 130.000 €. Berechnen Sie die Gewinnaufteilung nach §121 HGB.",
                    "context": "4% Zinsen auf Kapitalanteil, Rest nach Köpfen",
                    "table_structure": {
                        "headers": ["Name", "Kapitalanteil", "4% Zinsen", "Restgewinn", "Gesamtgewinn"],
                        "rows": [
                            {"label": "Bauer", "values": ["100.000 €", "", "", ""]},
                            {"label": "Müller", "values": ["150.000 €", "", "", ""]},
                        ],
                    },
                    "data": {"gewinn": 130000, "zinssatz": 0.04},
                    "external_url": "",
                    "solvable": True,
                    "skip_reason": "",
                },
                {
                    "id": "uebung_1",
                    "label": "Übung 1",
                    "task_type": "external_link",
                    "text": "Quiz auf lernnetz24.de",
                    "context": "",
                    "table_structure": None,
                    "data": {},
                    "external_url": "https://www.lernnetz24.de/bwl/hinweise/99.html",
                    "solvable": False,
                    "skip_reason": "Externer Link — Quiz muss im Browser bearbeitet werden",
                },
            ],
        }
    ),
    "Aufgaben-Löser für Schulaufgaben": json.dumps(
        {
            "answer": "Bauer erhält 25.000 € Gesamtgewinn, Müller erhält 66.000 €",
            "table_data": {
                "headers": ["Name", "Kapitalanteil", "4% Zinsen", "Restgewinn", "Gesamtgewinn"],
                "rows": [
                    {"label": "Bauer", "values": ["100.000 €", "4.000 €", "56.000 €", "60.000 €"]},
                    {"label": "Müller", "values": ["150.000 €", "6.000 €", "56.000 €", "62.000 €"]},
                    {"label": "Summe", "values": ["250.000 €", "10.000 €", "112.000 €", "122.000 €"]},
                ],
            },
            "calculation_steps": [
                "4% von 100.000 = 4.000 €",
                "4% von 150.000 = 6.000 €",
                "Restgewinn: 130.000 - 10.000 = 120.000 €",
                "Pro Kopf: 120.000 / 2 = 60.000 €",
            ],
            "confidence": 0.95,
        }
    ),
    # --- Document classification & template filling ---
    "Dokument-Klassifikator": json.dumps(
        {
            "documents": [
                {
                    "id": "doc_01",
                    "filename": "Projektantrag.docx",
                    "role": "template",
                    "reasoning": "Formular mit Leerfeldern",
                    "fields": [
                        {
                            "id": "field_01",
                            "label": "Projektbezeichnung",
                            "field_type": "text",
                            "location": "Zeile 3",
                            "current_value": "",
                            "max_length": 80,
                            "constraints": ["einzeilig"],
                        },
                        {
                            "id": "field_02",
                            "label": "Projektbeschreibung",
                            "field_type": "paragraph",
                            "location": "Abschnitt 2",
                            "current_value": "",
                            "max_length": 500,
                            "constraints": ["maximal halbe Seite"],
                        },
                    ],
                    "extracted_info": {},
                },
                {
                    "id": "doc_02",
                    "filename": "Anforderungen.txt",
                    "role": "source",
                    "reasoning": "Anforderungsliste",
                    "fields": None,
                    "extracted_info": {
                        "requirements": ["Webbasierte Lagerverwaltung", "Login-System", "REST API"],
                        "description": "Entwicklung einer Lagerverwaltung mit Flask",
                    },
                },
            ],
            "contradictions": [
                {
                    "topic": "Datenbankwahl",
                    "source_a": "Anforderungen.txt: SQLite",
                    "source_b": "Mündlich: PostgreSQL",
                    "recommendation": "SQLite (schriftlich dokumentiert)",
                }
            ],
        }
    ),
    "Template-Ausfüller": json.dumps(
        {
            "filled_fields": [
                {
                    "field_id": "field_01",
                    "label": "Projektbezeichnung",
                    "value": "Webbasierte Lagerverwaltung mit Flask und SQLite",
                    "fits_constraint": True,
                    "char_count": 48,
                },
                {
                    "field_id": "field_02",
                    "label": "Projektbeschreibung",
                    "value": "Entwicklung einer webbasierten Lagerverwaltung zur Verwaltung von Warenein- und -ausgängen.",
                    "fits_constraint": True,
                    "char_count": 91,
                },
            ],
            "warnings": ["Feld 'Zeitplan' hat nur 200 Zeichen"],
        }
    ),
    "Anforderungs-Auditor": json.dumps(
        {
            "findings": [
                {
                    "category": "gap",
                    "severity": "warning",
                    "title": "Abgabedatum nicht angegeben",
                    "detail": "Keines der Dokumente nennt ein Abgabedatum.",
                    "sources": ["Alle Dokumente"],
                    "quotes": [],
                    "recommendation": "Abgabedatum beim Lehrer erfragen",
                },
                {
                    "category": "ambiguity",
                    "severity": "info",
                    "title": "Datenbankwahl nicht eindeutig",
                    "detail": "Anforderungen erwähnen Datenbank ohne konkretes System.",
                    "sources": ["Anforderungen.txt"],
                    "quotes": ["muss Datenbank enthalten"],
                    "recommendation": "SQLite als Standard verwenden und dokumentieren",
                },
            ],
            "missing_information": [
                "Abgabedatum / Deadline",
                "Bewertungskriterien / Notenschlüssel",
                "Erlaubte Hilfsmittel",
            ],
            "completeness_assessment": {"score": 0.4, "reasoning": "Von 10 üblichen Pflichtangaben fehlen 6"},
        }
    ),
    "Anforderungs-Extraktor": json.dumps(
        {
            "requirements": [
                {
                    "id": "REQ-001",
                    "text": "Webbasierte Lagerverwaltung erstellen",
                    "source": "Anforderungen.txt",
                    "category": "functional",
                    "status": "clear",
                    "priority": "must",
                    "quote": "Entwicklung einer webbasierten Lagerverwaltung",
                },
                {
                    "id": "REQ-002",
                    "text": "Login-System implementieren",
                    "source": "Anforderungen.txt",
                    "category": "functional",
                    "status": "clear",
                    "priority": "must",
                    "quote": "Login-System",
                },
                {
                    "id": "REQ-003",
                    "text": "Datenbank verwenden",
                    "source": "Anforderungen.txt",
                    "category": "constraint",
                    "status": "ambiguous",
                    "priority": "must",
                    "quote": "muss Datenbank enthalten",
                },
            ],
            "implicit_requirements": [
                {
                    "id": "IMP-001",
                    "text": "Dokumentation in deutscher Sprache",
                    "reasoning": "Deutschsprachiger Unterricht",
                    "confidence": 0.95,
                }
            ],
            "requirement_count": {"total": 3, "clear": 2, "ambiguous": 1, "contradicted": 0, "gap": 0},
        }
    ),
    "Entscheidungs-Assistent": json.dumps(
        {
            "amendments": [
                {
                    "id": "AMD-001",
                    "resolves": "F-002",
                    "finding_title": "Widerspruch Datenbankwahl",
                    "decision": "SQLite wird verwendet",
                    "reasoning": "Schriftliche Vorgabe hat Vorrang. SQLite ist fuer den Projektumfang ausreichend.",
                    "source": "auto",
                    "alternatives_considered": ["PostgreSQL (muendlich empfohlen)", "MySQL (nicht erwaehnt)"],
                }
            ],
            "unresolvable": [{"finding_id": "F-003", "reason": "Abgabedatum muss beim Lehrer erfragt werden"}],
        }
    ),
    # --- Standard pipeline stages ---
    "Aufgaben-Parser": json.dumps(
        {
            "task_text": "Erstellen Sie eine Präsentation zum Thema IT-Sicherheit mit mindestens 8 Folien.",
            "subject": "IT-Sicherheit",
            "task_type": "presentation",
            "constraints": {
                "page_count": None,
                "slide_count": 8,
                "word_count": None,
                "language": "de",
                "format": "pptx",
                "due_date": None,
                "specific_requirements": ["mindestens 8 Folien"],
            },
            "raw_input_type": "text",
        }
    ),
    "Planungsassistent": json.dumps(
        {
            "title": "IT-Sicherheit im Unternehmen",
            "artifact_type": "pptx",
            "sections": [
                {
                    "id": "section_01",
                    "title": "IT-Sicherheit im Unternehmen",
                    "purpose": "Titelfolie",
                    "research_queries": [],
                    "estimated_length": "short",
                },
                {
                    "id": "section_02",
                    "title": "Was ist IT-Sicherheit?",
                    "purpose": "Definition",
                    "research_queries": ["IT-Sicherheit Definition"],
                    "estimated_length": "medium",
                },
                {
                    "id": "section_03",
                    "title": "Bedrohungen",
                    "purpose": "Angriffsvektoren",
                    "research_queries": ["Cyberangriffe"],
                    "estimated_length": "medium",
                },
                {
                    "id": "section_04",
                    "title": "Schutzmaßnahmen",
                    "purpose": "Maßnahmen",
                    "research_queries": ["IT-Sicherheit Maßnahmen"],
                    "estimated_length": "medium",
                },
                {
                    "id": "section_05",
                    "title": "Quellen",
                    "purpose": "Quellenangaben",
                    "research_queries": [],
                    "estimated_length": "short",
                },
            ],
            "style_notes": "Sachlich",
        }
    ),
    "Recherche-Assistent": json.dumps(
        {
            "sections": [
                {
                    "section_id": "section_01",
                    "findings": [{"content": "Titelfolie", "source": "llm_knowledge", "relevance": 1.0}],
                    "sufficient": True,
                },
                {
                    "section_id": "section_02",
                    "findings": [
                        {
                            "content": "IT-Sicherheit schützt Informationen und Systeme.",
                            "source": "llm_knowledge",
                            "relevance": 0.95,
                        },
                        {
                            "content": "CIA-Triad: Vertraulichkeit, Integrität, Verfügbarkeit.",
                            "source": "llm_knowledge",
                            "relevance": 0.95,
                        },
                    ],
                    "sufficient": True,
                },
                {
                    "section_id": "section_03",
                    "findings": [
                        {
                            "content": "Phishing ist die häufigste Angriffsmethode.",
                            "source": "llm_knowledge",
                            "relevance": 0.9,
                        },
                    ],
                    "sufficient": True,
                },
                {
                    "section_id": "section_04",
                    "findings": [
                        {
                            "content": "Firewalls, VPN, Verschlüsselung, Updates, Backups.",
                            "source": "llm_knowledge",
                            "relevance": 0.9,
                        },
                    ],
                    "sufficient": True,
                },
                {
                    "section_id": "section_05",
                    "findings": [{"content": "Quellenfolie", "source": "llm_knowledge", "relevance": 1.0}],
                    "sufficient": True,
                },
            ]
        }
    ),
    "Präsentations-Autor": json.dumps(
        {
            "title": "IT-Sicherheit im Unternehmen",
            "sections": [
                {
                    "section_id": "section_01",
                    "heading": "IT-Sicherheit im Unternehmen",
                    "content": "Grundlagen, Bedrohungen und Schutzmaßnahmen",
                    "bullet_points": [],
                    "speaker_notes": None,
                },
                {
                    "section_id": "section_02",
                    "heading": "Was ist IT-Sicherheit?",
                    "content": "Schutz von Informationen und IT-Systemen.",
                    "bullet_points": [
                        "CIA-Triad: Vertraulichkeit, Integrität, Verfügbarkeit",
                        "Gesetzliche Grundlage: BSI-Gesetz, DSGVO",
                        "Schutz vor unbefugtem Zugriff",
                    ],
                    "speaker_notes": "Die drei Schutzziele bilden die Basis.",
                },
                {
                    "section_id": "section_03",
                    "heading": "Bedrohungen",
                    "content": "Cyberangriffe nehmen zu.",
                    "bullet_points": ["Phishing — häufigste Methode", "Ransomware — Erpressung", "Social Engineering"],
                    "speaker_notes": "80% beginnen mit Phishing.",
                },
                {
                    "section_id": "section_04",
                    "heading": "Schutzmaßnahmen",
                    "content": "Technische und organisatorische Maßnahmen.",
                    "bullet_points": ["Firewalls und VPN", "Regelmäßige Updates und Backups", "Mitarbeiterschulungen"],
                    "speaker_notes": "Technik allein reicht nicht.",
                },
                {
                    "section_id": "section_05",
                    "heading": "Quellen",
                    "content": "",
                    "bullet_points": ["BSI", "OWASP Top 10"],
                    "speaker_notes": None,
                },
            ],
            "sources": ["BSI", "OWASP Top 10"],
        }
    ),
    "Dokument-Autor": json.dumps(
        {
            "title": "IT-Sicherheit",
            "sections": [
                {
                    "section_id": "section_01",
                    "heading": "Einleitung",
                    "content": "IT-Sicherheit ist ein zentrales Thema der modernen Informationstechnologie.",
                    "bullet_points": [],
                    "speaker_notes": None,
                },
            ],
            "sources": [],
        }
    ),
    "Aufgaben-Löser": json.dumps(
        {
            "title": "IT-Sicherheit Fragen",
            "sections": [
                {
                    "section_id": "section_01",
                    "heading": "Was ist IT-Sicherheit?",
                    "content": "IT-Sicherheit umfasst den Schutz von Informationen und IT-Systemen vor unbefugtem Zugriff.",
                    "bullet_points": [],
                    "speaker_notes": None,
                },
            ],
            "sources": [],
        }
    ),
    "Code-Generator": 'def main():\n    print("Hello World")\n\nif __name__ == "__main__":\n    main()\n',
}


class MockBackend:
    def __init__(self):
        self.name = "mock"
        self.is_free = True
        self.supports_vision = True
        self.max_context = 128_000
        self.model = "mock-model"
        self.calls = []

    async def complete(self, messages, temperature=0.3, max_tokens=4096, response_format=None):
        self.calls.append(messages)
        system_msg = messages[0]["content"] if messages else ""
        for key, response in DEFAULT_RESPONSES.items():
            if key in system_msg:
                return LLMResponse(content=response, model="mock", backend_name="mock")
        return LLMResponse(content="{}", model="mock", backend_name="mock")

    async def complete_vision(self, messages, temperature=0.3, max_tokens=4096):
        return await self.complete(messages, temperature, max_tokens)

    async def close(self):
        pass


def make_mock_router():
    """Create a mock router without needing real backends."""
    from schulpipeline.backends.router import BackendRouter

    config = PipelineConfig(
        backends={"mock": BackendConfig(name="mock", api_key="test", enabled=True)},
        cascade={
            stage: ["mock"]
            for stage in [
                "intake",
                "plan",
                "research",
                "synthesize",
                "artifact",
                "agent_codegen",
                "decompose",
                "solve",
                "classify_docs",
                "fill_template",
                "audit",
                "classify_report",
                "amendments",
            ]
        },
        research=ResearchConfig(enabled=False, use_web=False),
        output=OutputConfig(dir=tempfile.mkdtemp(), default_format="pptx", language="de"),
    )

    router = BackendRouter.__new__(BackendRouter)
    router.config = config
    router._backends = {"mock": MockBackend()}
    router._cooldowns = {}
    router._call_counts = {"mock": 0}
    router._total_cost = 0.0
    return config, router


# === Tests ===


class TestConfig(unittest.TestCase):
    def test_load_empty_config(self):
        config = load_config(path=None)
        self.assertIsInstance(config, PipelineConfig)
        self.assertEqual(config.output.language, "de")

    def test_backend_available(self):
        cfg = BackendConfig(name="groq", api_key="test-key", enabled=True)
        self.assertTrue(cfg.is_available)

    def test_backend_unavailable_no_key(self):
        cfg = BackendConfig(name="groq", api_key="", enabled=True)
        self.assertFalse(cfg.is_available)

    def test_backend_disabled(self):
        cfg = BackendConfig(name="groq", api_key="test", enabled=False)
        self.assertFalse(cfg.is_available)

    def test_cascade_filters(self):
        config = PipelineConfig(
            backends={"mock": BackendConfig(name="mock", api_key="test", enabled=True)},
            cascade={"plan": ["mock", "nonexistent"]},
        )
        self.assertEqual(config.cascade_for("plan"), ["mock"])

    def test_config_overrides(self):
        config = load_config(path=None, overrides={"log_level": "DEBUG"})
        self.assertEqual(config.log_level, "DEBUG")


class TestRouter(unittest.TestCase):
    def test_router_completes(self):
        config, router = make_mock_router()
        result = asyncio.run(router.complete(stage="plan", messages=[{"role": "user", "content": "test"}]))
        self.assertEqual(result.backend_name, "mock")

    def test_router_stats(self):
        _, router = make_mock_router()
        stats = router.stats()
        self.assertIn("backends_available", stats)
        self.assertIn("mock", stats["backends_available"])


class TestIntakeStage(unittest.TestCase):
    def test_text_intake(self):
        from schulpipeline.stages.intake import IntakeStage

        config, router = make_mock_router()
        stage = IntakeStage()
        result = asyncio.run(
            stage.run(
                {"raw_input": "Erstellen Sie eine Präsentation zum Thema IT-Sicherheit"},
                router,
                config,
            )
        )
        self.assertTrue(result.success, f"Intake failed: {result.errors}")
        self.assertEqual(result.data["task_type"], "presentation")
        self.assertEqual(result.data["constraints"]["format"], "pptx")


class TestPlanStage(unittest.TestCase):
    def test_plan_from_intake(self):
        from schulpipeline.stages.plan import PlanStage

        config, router = make_mock_router()
        context = {
            "intake": {
                "task_text": "Erstellen Sie eine Präsentation zum Thema IT-Sicherheit",
                "subject": "IT-Sicherheit",
                "task_type": "presentation",
                "constraints": {"language": "de", "format": "pptx", "slide_count": 8, "specific_requirements": []},
            }
        }
        stage = PlanStage()
        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"Plan failed: {result.errors}")
        self.assertIn("sections", result.data)
        self.assertGreater(len(result.data["sections"]), 0)


class TestArtifactBuilders(unittest.TestCase):
    def test_pptx_builder(self):
        from schulpipeline.artifacts.pptx_builder import build_pptx

        synthesis = json.loads(DEFAULT_RESPONSES["Präsentations-Autor"])
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = Path(f.name)
        try:
            build_pptx(synthesis, path)
            self.assertTrue(path.exists())
            self.assertGreater(path.stat().st_size, 1000)
        finally:
            path.unlink(missing_ok=True)

    def test_docx_builder(self):
        from schulpipeline.artifacts.docx_builder import build_docx

        synthesis = json.loads(DEFAULT_RESPONSES["Dokument-Autor"])
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            path = Path(f.name)
        try:
            build_docx(synthesis, path)
            self.assertTrue(path.exists())
            self.assertGreater(path.stat().st_size, 500)
        finally:
            path.unlink(missing_ok=True)

    def test_md_builder(self):
        from schulpipeline.artifacts.md_builder import build_md

        synthesis = json.loads(DEFAULT_RESPONSES["Aufgaben-Löser"])
        with tempfile.NamedTemporaryFile(suffix=".md", delete=False, mode="w") as f:
            path = Path(f.name)
        try:
            build_md(synthesis, path)
            self.assertTrue(path.exists())
            content = path.read_text()
            self.assertIn("IT-Sicherheit", content)
        finally:
            path.unlink(missing_ok=True)


class TestFullPipeline(unittest.TestCase):
    def test_text_to_pptx(self):
        from schulpipeline.pipeline import Pipeline

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)
        result = asyncio.run(pipeline.run("Erstellen Sie eine Präsentation zum Thema IT-Sicherheit mit 8 Folien."))

        self.assertTrue(result.success, f"Failed at {result.failed_stage}: {result.validation_errors}")
        self.assertIsNotNone(result.output_path)
        self.assertTrue(Path(result.output_path).exists())
        self.assertTrue(result.output_path.endswith(".pptx"))
        self.assertEqual(len(result.results), 5)

    def test_text_to_pptx_with_preset(self):
        """Full pipeline with a preset applied."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.presets import resolve_quick

        config, router = make_mock_router()
        preset = resolve_quick("fiae-praesi-itsec")
        pipeline = Pipeline(config, router)
        result = asyncio.run(pipeline.run("IT-Sicherheit im Unternehmen", preset=preset))

        self.assertTrue(result.success, f"Failed at {result.failed_stage}: {result.validation_errors}")
        self.assertTrue(result.output_path.endswith(".pptx"))

    def test_plan_only(self):
        from schulpipeline.pipeline import Pipeline

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)
        result = asyncio.run(pipeline.plan_only("Erstelle eine Präsentation zu Netzwerken"))

        self.assertTrue(result.success)
        self.assertEqual(len(result.results), 2)

    def test_pipeline_failure_reporting(self):
        from schulpipeline.pipeline import Pipeline

        config, router = make_mock_router()
        mock = router._backends["mock"]

        async def bad_complete(*args, **kwargs):
            return LLMResponse(content="not json!!!", model="mock", backend_name="mock")

        mock.complete = bad_complete

        pipeline = Pipeline(config, router)
        result = asyncio.run(pipeline.run("test"))

        self.assertFalse(result.success)
        self.assertEqual(result.failed_stage, "intake")


class TestPresets(unittest.TestCase):
    def test_resolve_quick_preset(self):
        from schulpipeline.presets import resolve_quick

        preset = resolve_quick("fiae-praesi-itsec")
        self.assertEqual(preset.output_format, "pptx")
        self.assertEqual(preset.subject_key, "it_sicherheit")
        self.assertEqual(preset.section_count, 10)
        self.assertIn("CIA-Triad", preset.vocabulary_hints)
        self.assertIn("Berufsschul-Niveau", preset.system_context)

    def test_resolve_manual_preset(self):
        from schulpipeline.presets import resolve_preset

        preset = resolve_preset("ausarbeitung", "wirtschaft")
        self.assertEqual(preset.output_format, "docx")
        self.assertEqual(preset.subject_key, "wirtschaft")
        self.assertEqual(preset.style, "prose")
        self.assertIn("formell", preset.quality_instructions.lower())

    def test_resolve_with_overrides(self):
        from schulpipeline.presets import resolve_quick

        preset = resolve_quick("fiae-praesi-itsec", overrides={"section_count": 15})
        self.assertEqual(preset.section_count, 15)

    def test_unknown_preset_raises(self):
        from schulpipeline.presets import resolve_quick

        with self.assertRaises(ValueError):
            resolve_quick("nonexistent-preset")

    def test_unknown_output_type_raises(self):
        from schulpipeline.presets import resolve_preset

        with self.assertRaises(ValueError):
            resolve_preset("nonexistent", "it_sicherheit")

    def test_unknown_subject_raises(self):
        from schulpipeline.presets import resolve_preset

        with self.assertRaises(ValueError):
            resolve_preset("praesentation", "nonexistent")

    def test_system_context_contains_domain(self):
        from schulpipeline.presets import resolve_quick

        preset = resolve_quick("fiae-praesi-netzwerk")
        self.assertIn("OSI-Modell", preset.system_context)
        self.assertIn("TCP/IP", preset.system_context)

    def test_quality_instructions_for_presentation(self):
        from schulpipeline.presets import resolve_quick

        preset = resolve_quick("fiae-praesi-itsec")
        self.assertIn("Stichpunkte", preset.quality_instructions)
        self.assertIn("Speaker Notes", preset.quality_instructions)

    def test_quality_instructions_for_aufgaben(self):
        from schulpipeline.presets import resolve_quick

        preset = resolve_quick("fiae-aufgaben-prog")
        self.assertIn("Direkte Antworten", preset.quality_instructions)

    def test_list_presets(self):
        from schulpipeline.presets import list_presets

        data = list_presets()
        self.assertIn("output_types", data)
        self.assertIn("subjects", data)
        self.assertIn("quick", data)
        self.assertIn("praesentation", data["output_types"])
        self.assertIn("it_sicherheit", data["subjects"])

    def test_all_quick_presets_resolve(self):
        """Every quick preset must resolve without errors."""
        from schulpipeline.presets import QUICK_PRESETS, resolve_quick

        for key in QUICK_PRESETS:
            preset = resolve_quick(key)
            self.assertIsNotNone(preset.system_context)
            self.assertIn(
                preset.output_format,
                ("pptx", "docx", "md", "project", "worksheet", "template_fill", "audit", "requirements_report"),
            )

    def test_english_subject_sets_language(self):
        from schulpipeline.presets import resolve_preset

        preset = resolve_preset("aufgaben", "englisch")
        self.assertEqual(preset.language, "en")


class TestSession(unittest.TestCase):
    def setUp(self):
        self.tmpdir = tempfile.mkdtemp()

    def tearDown(self):
        import shutil

        shutil.rmtree(self.tmpdir, ignore_errors=True)

    def test_create_and_load_session(self):
        from schulpipeline.session import SessionStore

        store = SessionStore(sessions_dir=os.path.join(self.tmpdir, "sessions"))

        session = store.create(
            task_input="Test task",
            preset_key="fiae-praesi-itsec",
            subject="it_sicherheit",
            tags=["test"],
        )
        self.assertEqual(session.status, "created")
        self.assertEqual(len(session.id), 8)

        loaded = store.load(session.id)
        self.assertIsNotNone(loaded)
        self.assertEqual(loaded.task_input, "Test task")
        self.assertEqual(loaded.preset_key, "fiae-praesi-itsec")
        self.assertEqual(loaded.tags, ["test"])

    def test_session_list_and_filter(self):
        from schulpipeline.session import SessionStore

        store = SessionStore(sessions_dir=os.path.join(self.tmpdir, "sessions"))

        s1 = store.create(task_input="Task 1", subject="it_sicherheit")
        store.create(task_input="Task 2", subject="wirtschaft")
        s1.status = "completed"
        store.save(s1)

        all_sessions = store.list_sessions()
        self.assertEqual(len(all_sessions), 2)

        completed = store.list_sessions(status="completed")
        self.assertEqual(len(completed), 1)
        self.assertEqual(completed[0]["id"], s1.id)

        wirtschaft = store.list_sessions(subject="wirtschaft")
        self.assertEqual(len(wirtschaft), 1)

    def test_session_delete(self):
        from schulpipeline.session import SessionStore

        store = SessionStore(sessions_dir=os.path.join(self.tmpdir, "sessions"))

        session = store.create(task_input="Delete me")
        self.assertTrue(store.delete(session.id))
        self.assertIsNone(store.load(session.id))
        self.assertFalse(store.delete("nonexistent"))

    def test_session_display_title(self):
        from schulpipeline.session import Session, StageSnapshot

        session = Session(
            id="test1234",
            created_at="",
            updated_at="",
            task_input="A very long task that should be truncated to sixty characters at most ok",
            input_type="text",
        )
        self.assertEqual(len(session.display_title), 60)

        # With plan data, title comes from plan
        session.completed_stages.append(
            StageSnapshot(
                name="plan",
                success=True,
                data={"title": "IT-Sicherheit im Unternehmen"},
                errors=[],
                elapsed_ms=0,
                backend_used="",
                completed_at="",
            )
        )
        self.assertEqual(session.display_title, "IT-Sicherheit im Unternehmen")

    def test_session_stage_data(self):
        from schulpipeline.session import Session, StageSnapshot

        session = Session(id="t", created_at="", updated_at="", task_input="x", input_type="text")
        session.completed_stages.append(
            StageSnapshot(
                name="intake",
                success=True,
                data={"task_text": "hello"},
                errors=[],
                elapsed_ms=100,
                backend_used="mock",
                completed_at="",
            )
        )
        self.assertEqual(session.stage_data["intake"]["task_text"], "hello")

    def test_session_is_resumable(self):
        from schulpipeline.session import Session

        s = Session(id="t", created_at="", updated_at="", task_input="x", input_type="text")
        s.status = "completed"
        self.assertFalse(s.is_resumable)
        s.status = "failed"
        self.assertTrue(s.is_resumable)
        s.status = "paused"
        self.assertTrue(s.is_resumable)

    def test_session_runner_full(self):
        """Full pipeline through session runner."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.session import SessionRunner, SessionStore

        config, router = make_mock_router()
        config.output.dir = self.tmpdir
        store = SessionStore(sessions_dir=os.path.join(self.tmpdir, "sessions"))
        pipeline = Pipeline(config, router)
        runner = SessionRunner(store, pipeline, router)

        session = store.create(task_input="IT-Sicherheit Präsentation")
        session = asyncio.run(runner.run(session))

        self.assertEqual(session.status, "completed")
        self.assertIsNotNone(session.output_path)
        self.assertEqual(len(session.completed_stages), 5)

        # Verify persistence
        loaded = store.load(session.id)
        self.assertEqual(loaded.status, "completed")
        self.assertEqual(len(loaded.completed_stages), 5)

    def test_session_runner_resume(self):
        """Resume from a failed session after fixing the issue."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.session import SessionRunner, SessionStore, StageSnapshot

        config, router = make_mock_router()
        config.output.dir = self.tmpdir
        store = SessionStore(sessions_dir=os.path.join(self.tmpdir, "sessions"))
        pipeline = Pipeline(config, router)
        runner = SessionRunner(store, pipeline, router)

        # Create a session that already has intake + plan done
        session = store.create(task_input="Test resume")
        session.completed_stages = [
            StageSnapshot(
                name="intake",
                success=True,
                data=json.loads(DEFAULT_RESPONSES["Aufgaben-Parser"]),
                errors=[],
                elapsed_ms=100,
                backend_used="mock",
                completed_at="",
            ),
            StageSnapshot(
                name="plan",
                success=True,
                data=json.loads(DEFAULT_RESPONSES["Planungsassistent"]),
                errors=[],
                elapsed_ms=100,
                backend_used="mock",
                completed_at="",
            ),
        ]
        session.status = "failed"
        session.failed_stage = "research"
        store.save(session)

        # Resume — should skip intake and plan
        session = asyncio.run(runner.run(session))
        self.assertEqual(session.status, "completed")
        # Should have 5 stages total (2 old + 3 new)
        self.assertEqual(len(session.completed_stages), 5)

    def test_session_retry_from_stage(self):
        """Retry from a specific stage, dropping later stages."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.session import SessionRunner, SessionStore, StageSnapshot

        config, router = make_mock_router()
        config.output.dir = self.tmpdir
        store = SessionStore(sessions_dir=os.path.join(self.tmpdir, "sessions"))
        pipeline = Pipeline(config, router)
        runner = SessionRunner(store, pipeline, router)

        # Session with all stages done
        session = store.create(task_input="Test retry")
        session.completed_stages = [
            StageSnapshot(
                name="intake",
                success=True,
                data=json.loads(DEFAULT_RESPONSES["Aufgaben-Parser"]),
                errors=[],
                elapsed_ms=100,
                backend_used="mock",
                completed_at="",
            ),
            StageSnapshot(
                name="plan",
                success=True,
                data=json.loads(DEFAULT_RESPONSES["Planungsassistent"]),
                errors=[],
                elapsed_ms=100,
                backend_used="mock",
                completed_at="",
            ),
            StageSnapshot(
                name="research",
                success=True,
                data=json.loads(DEFAULT_RESPONSES["Recherche-Assistent"]),
                errors=[],
                elapsed_ms=100,
                backend_used="mock",
                completed_at="",
            ),
        ]
        session.status = "completed"
        store.save(session)

        # Retry from synthesize — drops research results, re-runs from synthesize
        session = asyncio.run(runner.retry_from(session, "research"))
        self.assertEqual(session.status, "completed")
        # intake + plan kept, research + synthesize + artifact re-run
        self.assertEqual(len(session.completed_stages), 5)


class TestAgents(unittest.TestCase):
    def setUp(self):
        self.tmpdir = tempfile.mkdtemp()

    def tearDown(self):
        import shutil

        shutil.rmtree(self.tmpdir, ignore_errors=True)

    def test_build_project_spec(self):
        from schulpipeline.agents import build_project_spec

        synthesis = {
            "title": "Lagerverwaltung",
            "sections": [
                {
                    "section_id": "s1",
                    "heading": "Datenbankmodell",
                    "content": "SQLite Datenbank mit Tabellen für Produkte und Bestellungen",
                    "bullet_points": ["create_product()", "get_all_products()"],
                    "speaker_notes": None,
                },
                {
                    "section_id": "s2",
                    "heading": "REST API",
                    "content": "Flask-basierte API mit CRUD Endpunkten",
                    "bullet_points": ["GET /products", "POST /products"],
                    "speaker_notes": None,
                },
            ],
            "sources": [],
        }
        intake = {
            "task_text": "Erstellen Sie eine Lagerverwaltung mit Python und Flask",
            "constraints": {"specific_requirements": ["REST API", "SQLite"]},
        }

        spec = build_project_spec(synthesis, intake)
        self.assertEqual(spec.language, "python")
        self.assertEqual(spec.framework, "flask")
        self.assertEqual(len(spec.modules), 2)
        self.assertIn("flask", spec.dependencies)

    def test_project_spec_to_prompt(self):
        from schulpipeline.agents import FileSpec, ModuleSpec, ProjectSpec

        spec = ProjectSpec(
            title="Test Project",
            description="A test",
            language="python",
            framework="flask",
            modules=[
                ModuleSpec(
                    name="API",
                    purpose="REST endpoints",
                    files=[
                        FileSpec(path="src/api.py", description="Endpoints", key_functions=["get_items"]),
                    ],
                ),
            ],
            dependencies=["flask"],
            requirements=["include tests"],
        )

        prompt = spec.to_prompt()
        self.assertIn("Test Project", prompt)
        self.assertIn("flask", prompt)
        self.assertIn("get_items", prompt)
        self.assertIn("include tests", prompt)

    def test_local_llm_agent(self):
        """Local LLM agent generates files."""
        from schulpipeline.agents import FileSpec, LocalLLMAgent, ModuleSpec, ProjectSpec

        _, router = make_mock_router()
        # Override cascade to include agent_codegen stage
        router.config.cascade["agent_codegen"] = ["mock"]

        agent = LocalLLMAgent(router)
        self.assertTrue(agent.is_free)

        spec = ProjectSpec(
            title="Test",
            description="A test project",
            language="python",
            modules=[
                ModuleSpec(
                    name="main",
                    purpose="Entry point",
                    files=[
                        FileSpec(path="src/main.py", description="Main module", key_functions=["main"]),
                    ],
                ),
            ],
        )

        cost = asyncio.run(agent.estimate_cost(spec))
        self.assertEqual(cost, 0.0)

        output_dir = Path(self.tmpdir) / "test_project"
        result = asyncio.run(agent.execute(spec, output_dir))

        self.assertTrue(result.success, f"Agent failed: {result.errors}")
        self.assertTrue(any("main.py" in f for f in result.files_created))
        self.assertTrue(any("README" in f for f in result.files_created))

    def test_detect_language(self):
        from schulpipeline.agents import _detect_language

        self.assertEqual(_detect_language("erstelle eine flask webapp"), "python")
        self.assertEqual(_detect_language("react single page application"), "javascript")
        self.assertEqual(_detect_language("spring boot microservice"), "java")
        self.assertEqual(_detect_language("ein kleines tool"), "python")  # default


class TestWorksheet(unittest.TestCase):
    """Tests for worksheet decomposition, solving, and formatting."""

    def test_decompose_stage(self):
        """DecomposeStage parses messy input into structured tasks."""
        from schulpipeline.worksheet import DecomposeStage

        config, router = make_mock_router()
        stage = DecomposeStage()
        self.assertEqual(stage.name, "decompose")

        context = {
            "intake": {
                "task_text": SAMPLE_WORKSHEET_INPUT,
                "subject": "Wirtschaft",
                "task_type": "worksheet",
            },
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"Decompose failed: {result.errors}")
        self.assertEqual(result.data["title"], "Gewinnverteilung OHG/KG")
        self.assertEqual(len(result.data["tasks"]), 2)

        # First task should be solvable
        t0 = result.data["tasks"][0]
        self.assertTrue(t0["solvable"])
        self.assertEqual(t0["task_type"], "table_fill")
        self.assertIsNotNone(t0.get("table_structure"))

        # Second task should be unsolvable (external link)
        t1 = result.data["tasks"][1]
        self.assertFalse(t1["solvable"])
        self.assertEqual(t1["task_type"], "external_link")

    def test_solve_stage(self):
        """SolveStage solves decomposed tasks and separates unsolvable ones."""
        from schulpipeline.worksheet import SolveStage

        config, router = make_mock_router()
        stage = SolveStage()

        decomposed = json.loads(DEFAULT_RESPONSES["Aufgaben-Parser für Schularbeitsblätter"])
        context = {
            "intake": {"task_text": "test", "subject": "Wirtschaft", "task_type": "worksheet"},
            "decompose": decomposed,
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"Solve failed: {result.errors}")

        # Should have 1 solved task (aufgabe_1a) and 1 unsolvable (uebung_1)
        self.assertEqual(len(result.data["solved_tasks"]), 1)
        self.assertEqual(len(result.data["unsolvable_tasks"]), 1)

        solved = result.data["solved_tasks"][0]
        self.assertIn("table_data", solved["solution"])
        self.assertIn("calculation_steps", solved["solution"])
        self.assertEqual(len(solved["solution"]["calculation_steps"]), 4)

    def test_format_worksheet_as_md(self):
        """Markdown formatter produces correct structure."""
        from schulpipeline.worksheet import format_worksheet_as_md

        solve_result = {
            "title": "Gewinnverteilung",
            "subject": "Wirtschaft",
            "solved_tasks": [
                {
                    "task": {"label": "Aufgabe 1a", "text": "Berechnen Sie..."},
                    "solution": {
                        "answer": "Bauer: 60.000 €, Müller: 62.000 €",
                        "table_data": {
                            "headers": ["Name", "Kapital", "Zinsen", "Rest", "Gesamt"],
                            "rows": [
                                {"label": "Bauer", "values": ["100.000 €", "4.000 €", "56.000 €", "60.000 €"]},
                                {"label": "Müller", "values": ["150.000 €", "6.000 €", "56.000 €", "62.000 €"]},
                            ],
                        },
                        "calculation_steps": ["4% von 100.000 = 4.000 €"],
                        "confidence": 0.95,
                    },
                }
            ],
            "unsolvable_tasks": [
                {"label": "Übung 1", "skip_reason": "Externer Link", "external_url": "https://example.com"}
            ],
        }

        md = format_worksheet_as_md(solve_result)
        self.assertIn("# Gewinnverteilung", md)
        self.assertIn("Wirtschaft", md)
        self.assertIn("| Bauer |", md)
        self.assertIn("4.000 €", md)
        self.assertIn("Rechenweg", md)
        self.assertIn("Nicht bearbeitbare Aufgaben", md)
        self.assertIn("Externer Link", md)

    def test_format_worksheet_as_docx(self):
        """DOCX formatter produces a valid file with tables."""
        from schulpipeline.worksheet import format_worksheet_as_docx

        solve_result = {
            "title": "Testblatt",
            "subject": "BWL",
            "solved_tasks": [
                {
                    "task": {"label": "Aufgabe 1", "text": "Berechnen Sie den Gewinn."},
                    "solution": {
                        "answer": "Ergebnis: 100.000 €",
                        "table_data": {
                            "headers": ["Name", "Betrag"],
                            "rows": [
                                {"label": "A", "values": ["50.000 €"]},
                                {"label": "B", "values": ["50.000 €"]},
                            ],
                        },
                        "calculation_steps": ["50.000 + 50.000 = 100.000"],
                        "confidence": 1.0,
                    },
                }
            ],
            "unsolvable_tasks": [],
        }

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            path = f.name

        try:
            format_worksheet_as_docx(solve_result, path)
            self.assertTrue(Path(path).exists())
            self.assertGreater(Path(path).stat().st_size, 0)

            # Verify it's a valid docx
            from docx import Document

            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
            self.assertIn("Testblatt", text)
            self.assertIn("Aufgabe 1", text)
            self.assertIn("100.000 €", text)

            # Verify table exists
            self.assertGreaterEqual(len(doc.tables), 1)
            table = doc.tables[0]
            self.assertEqual(table.rows[0].cells[0].text, "Name")
            self.assertEqual(table.rows[1].cells[0].text, "A")
        finally:
            Path(path).unlink(missing_ok=True)

    def test_worksheet_preset_resolves(self):
        """Worksheet presets resolve correctly."""
        from schulpipeline.presets import resolve_preset

        preset = resolve_preset("arbeitsblatt", "wirtschaft")
        self.assertEqual(preset.output_format, "worksheet")
        self.assertTrue(preset.output_constraints.get("worksheet_mode"))
        self.assertTrue(preset.output_constraints.get("show_calculation_steps"))
        self.assertTrue(preset.output_constraints.get("german_number_format"))

    def test_worksheet_quick_presets(self):
        """All FIAE-Blatt quick presets resolve."""
        from schulpipeline.presets import resolve_quick

        for key in ["fiae-blatt-wirtschaft", "fiae-blatt-itsec", "fiae-blatt-prog"]:
            preset = resolve_quick(key)
            self.assertIsNotNone(preset, f"Quick preset {key} should resolve")
            self.assertEqual(preset.output_format, "worksheet")

    def test_pipeline_selects_worksheet_stages(self):
        """Pipeline uses decompose→solve flow for worksheet presets."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.presets import resolve_preset

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)

        preset = resolve_preset("arbeitsblatt", "wirtschaft")
        context = {"preset": preset}

        stages = pipeline._select_stages(context)
        stage_names = [s.name for s in stages]
        self.assertEqual(stage_names, ["intake", "decompose", "solve"])

    def test_pipeline_standard_stages_for_presentation(self):
        """Pipeline still uses standard stages for non-worksheet presets."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.presets import resolve_preset

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)

        preset = resolve_preset("praesentation", "it_sicherheit")
        context = {"preset": preset}

        stages = pipeline._select_stages(context)
        stage_names = [s.name for s in stages]
        self.assertEqual(stage_names, ["intake", "plan", "research", "synthesize", "artifact"])


# --- Sample real-world worksheet input for testing ---
SAMPLE_WORKSHEET_INPUT = """Übung 1:
😎Testen Sie sich: Quiz (Klick auf "Sofort üben)
https://www.lernnetz24.de/bwl/hinweise/99.html

Übung 2: Gewinnverteilung

Aufgabe 1 a) Der Gewinn der Baumeister OHG beträgt 130.000,00 €. Die Einlage von Herrn Bauer beläuft sich auf 100.000,00 €, die von Frau Müller auf 150.000,00 €. Die Verteilung soll nach altem Recht § 121 HGB erfolgen.
"Von dem Jahresgewinne gebührt jedem Gesellschafter zunächst ein Anteil in Höhe von vier vom Hundert seines Kapitalanteils."
Berechnen Sie die Gewinnaufteilung!
Name Kapitalanteil 4% Zinsen Restgewinn Gesamtgewinn
Bauer
Müller
Summe
"""


class TestDocuments(unittest.TestCase):
    """Tests for document role classification and template filling."""

    def test_classify_docs_stage(self):
        """ClassifyDocsStage correctly classifies documents by role."""
        from schulpipeline.documents import ClassifyDocsStage

        config, router = make_mock_router()
        stage = ClassifyDocsStage()
        self.assertEqual(stage.name, "classify_docs")

        context = {
            "documents": [
                {
                    "filename": "Projektantrag.docx",
                    "content_type": "docx",
                    "content": "Projektbezeichnung: ___\nProjektbeschreibung: ___\n",
                },
                {
                    "filename": "Anforderungen.txt",
                    "content_type": "text",
                    "content": "Webbasierte Lagerverwaltung\n- Login-System\n- REST API\n",
                },
            ],
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"Classify failed: {result.errors}")

        docs = result.data["documents"]
        self.assertEqual(len(docs), 2)

        template_doc = next(d for d in docs if d["role"] == "template")
        self.assertEqual(template_doc["filename"], "Projektantrag.docx")
        self.assertIsNotNone(template_doc.get("fields"))
        self.assertGreaterEqual(len(template_doc["fields"]), 1)

        source_doc = next(d for d in docs if d["role"] == "source")
        self.assertEqual(source_doc["filename"], "Anforderungen.txt")

    def test_fill_template_stage(self):
        """FillTemplateStage generates field values from source info."""
        from schulpipeline.documents import FillTemplateStage

        config, router = make_mock_router()
        stage = FillTemplateStage()

        classify_result = json.loads(DEFAULT_RESPONSES["Dokument-Klassifikator"])
        context = {
            "intake": {"task_text": "Projektantrag ausfüllen", "subject": "Programmierung"},
            "classify_docs": classify_result,
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"Fill failed: {result.errors}")

        filled = result.data["filled_templates"]
        self.assertEqual(len(filled), 1)
        self.assertGreaterEqual(len(filled[0]["fields_filled"]), 1)

        # Check that field values respect char limits
        for field in filled[0]["fields_filled"]:
            self.assertIn("value", field)
            self.assertTrue(len(field["value"]) > 0)

    def test_apply_to_docx(self):
        """Template application fills DOCX fields while preserving layout."""
        from docx import Document

        from schulpipeline.documents import apply_to_docx

        # Create a simple template
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            template_path = f.name

        doc = Document()
        doc.add_heading("Projektantrag", level=0)
        doc.add_paragraph("Projektbezeichnung: {{PROJEKT_NAME}}")
        doc.add_paragraph("Beschreibung: {{PROJEKT_DESC}}")

        # Add a table with a placeholder
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Feld"
        table.rows[0].cells[1].text = "Wert"
        table.rows[1].cells[0].text = "Sprache"
        table.rows[1].cells[1].text = "{{SPRACHE}}"
        doc.save(template_path)

        output_path = template_path.replace(".docx", "_filled.docx")

        try:
            filled_fields = [
                {"label": "{{PROJEKT_NAME}}", "value": "Lagerverwaltung"},
                {"label": "{{PROJEKT_DESC}}", "value": "Webbasierte Lagerverwaltung mit Flask"},
                {"label": "{{SPRACHE}}", "value": "Python"},
            ]

            apply_to_docx(template_path, filled_fields, output_path)

            # Verify output exists and is valid
            self.assertTrue(Path(output_path).exists())
            result_doc = Document(output_path)

            full_text = "\n".join(p.text for p in result_doc.paragraphs)
            self.assertIn("Lagerverwaltung", full_text)
            self.assertIn("Webbasierte Lagerverwaltung mit Flask", full_text)
            self.assertNotIn("{{PROJEKT_NAME}}", full_text)
            self.assertNotIn("{{PROJEKT_DESC}}", full_text)

            # Check table was filled
            table_text = result_doc.tables[0].rows[1].cells[1].text
            self.assertEqual(table_text, "Python")
        finally:
            Path(template_path).unlink(missing_ok=True)
            Path(output_path).unlink(missing_ok=True)

    def test_apply_to_pptx(self):
        """Template application fills PPTX placeholders."""
        from pptx import Presentation
        from pptx.util import Inches

        from schulpipeline.documents import apply_to_pptx

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            template_path = f.name

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        txBox.text_frame.text = "Projekttitel: {{TITLE}}"
        prs.save(template_path)

        output_path = template_path.replace(".pptx", "_filled.pptx")

        try:
            filled_fields = [
                {"label": "{{TITLE}}", "value": "Lagerverwaltung v2.0"},
            ]

            apply_to_pptx(template_path, filled_fields, output_path)

            result_prs = Presentation(output_path)
            slide_text = ""
            for slide in result_prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text += shape.text_frame.text
            self.assertIn("Lagerverwaltung v2.0", slide_text)
            self.assertNotIn("{{TITLE}}", slide_text)
        finally:
            Path(template_path).unlink(missing_ok=True)
            Path(output_path).unlink(missing_ok=True)

    def test_template_presets(self):
        """Template presets resolve correctly."""
        from schulpipeline.presets import resolve_preset, resolve_quick

        preset = resolve_preset("vorlage", "programmierung")
        self.assertEqual(preset.output_format, "template_fill")
        self.assertTrue(preset.output_constraints.get("template_mode"))

        antrag = resolve_preset("projektantrag", "programmierung")
        self.assertTrue(antrag.output_constraints.get("max_pages") == 1)

        quick = resolve_quick("fiae-projektantrag")
        self.assertEqual(quick.output_format, "template_fill")

    def test_pipeline_selects_template_stages(self):
        """Pipeline uses classify→fill flow for template presets."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.presets import resolve_preset

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)

        preset = resolve_preset("vorlage", "programmierung")
        context = {"preset": preset}

        stages = pipeline._select_stages(context)
        stage_names = [s.name for s in stages]
        self.assertEqual(stage_names, ["intake", "classify_docs", "audit", "fill_template"])

    def test_contradiction_detection(self):
        """Classifier detects contradictions between documents."""
        classify_result = json.loads(DEFAULT_RESPONSES["Dokument-Klassifikator"])
        contradictions = classify_result.get("contradictions", [])
        self.assertEqual(len(contradictions), 1)
        self.assertEqual(contradictions[0]["topic"], "Datenbankwahl")
        self.assertIn("SQLite", contradictions[0]["recommendation"])

    def test_fill_respects_char_limit(self):
        """Filled fields report character counts for constraint checking."""
        fill_result = json.loads(DEFAULT_RESPONSES["Template-Ausfüller"])
        for field in fill_result["filled_fields"]:
            self.assertIn("char_count", field)
            self.assertIn("fits_constraint", field)
            # Verify reported char_count matches actual value length
            self.assertEqual(field["char_count"], len(field["value"]))


class TestAudit(unittest.TestCase):
    """Tests for specification audit — the most valuable pipeline output."""

    def test_deterministic_field_feasibility(self):
        """Detect when template fields can't hold required content."""
        from schulpipeline.audit import check_template_field_feasibility

        template = {
            "filename": "Projektantrag.docx",
            "fields": [
                {"label": "Name", "field_type": "text", "max_length": 10},
                {"label": "Beschreibung", "field_type": "paragraph", "max_length": 50},
            ],
        }
        requirements = [
            "Webbasierte Lagerverwaltung",
            "Login-System mit Rollenverwaltung",
            "REST API",
            "Datenbankanbindung",
            "Automatisierte Tests",
        ]

        findings = check_template_field_feasibility(template, requirements)
        # Should flag the tiny Name field
        tiny_fields = [f for f in findings if "10 Zeichen" in f.get("detail", "")]
        self.assertGreaterEqual(len(tiny_fields), 1)

    def test_deterministic_page_constraint(self):
        """Detect when content doesn't fit page limits."""
        from schulpipeline.audit import check_page_constraint

        # Template with many paragraph fields vs 1 page limit
        template = {
            "filename": "Antrag.docx",
            "fields": [
                {"label": f"Section {i}", "field_type": "paragraph", "max_length": 500}
                for i in range(15)  # 15 paragraphs = way too much for 1 page
            ],
        }

        findings = check_page_constraint(template, max_pages=1, requirements=[])
        blockers = [f for f in findings if f.get("severity") == "blocker"]
        self.assertGreaterEqual(len(blockers), 1)
        self.assertIn("passt nicht", blockers[0]["title"])

    def test_deterministic_page_constraint_passes(self):
        """No findings when content fits page limit."""
        from schulpipeline.audit import check_page_constraint

        template = {
            "filename": "Kurz.docx",
            "fields": [
                {"label": "Titel", "field_type": "text", "max_length": 80},
                {"label": "Beschreibung", "field_type": "paragraph", "max_length": 200},
            ],
        }

        findings = check_page_constraint(template, max_pages=1, requirements=[])
        self.assertEqual(len(findings), 0)

    def test_deterministic_contradiction_detection(self):
        """Detect contradictions between documents."""
        from schulpipeline.audit import check_contradictions_deterministic

        documents = [
            {
                "filename": "Anforderungen.docx",
                "extracted_info": {"datenbank": "SQLite", "framework": "Flask"},
            },
            {
                "filename": "Mündlich.txt",
                "extracted_info": {"datenbank": "PostgreSQL", "framework": "Flask"},
            },
        ]

        findings = check_contradictions_deterministic(documents)
        contradictions = [f for f in findings if f["category"] == "contradiction"]
        self.assertEqual(len(contradictions), 1)
        self.assertIn("datenbank", contradictions[0]["title"].lower())

    def test_deterministic_no_false_contradictions(self):
        """No contradictions when documents agree."""
        from schulpipeline.audit import check_contradictions_deterministic

        documents = [
            {"filename": "A.docx", "extracted_info": {"framework": "Flask"}},
            {"filename": "B.docx", "extracted_info": {"framework": "Flask"}},
        ]

        findings = check_contradictions_deterministic(documents)
        self.assertEqual(len(findings), 0)

    def test_missing_references(self):
        """Detect references to documents not provided."""
        from schulpipeline.audit import check_missing_references

        documents = [
            {
                "filename": "Aufgabe.docx",
                "role": "source",
                "content": "Bearbeiten Sie die Aufgabe laut Bewertungsschema. Siehe auch Anlage 3.",
            },
        ]

        findings = check_missing_references(documents)
        references = [f for f in findings if f["category"] == "underdefined"]
        # Should find Bewertungsschema and/or Anlage
        self.assertGreaterEqual(len(references), 1)

    def test_audit_stage_full(self):
        """Full audit stage combines deterministic + LLM findings."""
        from schulpipeline.audit import AuditStage

        config, router = make_mock_router()
        stage = AuditStage()

        classify_result = json.loads(DEFAULT_RESPONSES["Dokument-Klassifikator"])
        # Add content to documents for reference checking
        for doc in classify_result["documents"]:
            doc["content"] = doc.get("content", f"Inhalt von {doc['filename']}")

        context = {
            "classify_docs": classify_result,
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"Audit failed: {result.errors}")

        # Should have findings (at minimum the contradiction from classify_docs)
        self.assertGreater(result.data["summary"]["total_findings"], 0)
        self.assertIn("verdict", result.data["summary"])
        self.assertIsInstance(result.data["summary"]["completeness_score"], float)
        self.assertIsInstance(result.data["summary"]["feasibility_score"], float)

        # All findings should have required fields
        for f in result.data["findings"]:
            self.assertIn("id", f)
            self.assertIn("category", f)
            self.assertIn("severity", f)
            self.assertIn("title", f)

    def test_audit_report_md(self):
        """Markdown report is readable and structured."""
        from schulpipeline.audit import format_audit_as_md

        audit_data = {
            "title": "Vorgaben-Audit: Lagerverwaltung",
            "documents_analyzed": ["Projektantrag.docx", "Anforderungen.txt"],
            "findings": [
                {
                    "id": "F-001",
                    "category": "contradiction",
                    "severity": "blocker",
                    "title": "Widersprüchliche DB-Vorgabe",
                    "detail": "Anforderungen sagen SQLite, mündlich PostgreSQL.",
                    "sources": ["Anforderungen.txt", "Mündlich"],
                    "quotes": ["SQLite verwenden", "Nehmt PostgreSQL"],
                    "recommendation": "Schriftliche Vorgabe (SQLite) verwenden.",
                },
                {
                    "id": "F-002",
                    "category": "gap",
                    "severity": "warning",
                    "title": "Abgabedatum fehlt",
                    "detail": "Kein Dokument nennt ein Abgabedatum.",
                    "sources": ["Alle"],
                    "quotes": [],
                    "recommendation": "Abgabedatum erfragen.",
                },
            ],
            "summary": {
                "total_findings": 2,
                "blockers": 1,
                "warnings": 1,
                "info": 0,
                "completeness_score": 0.4,
                "feasibility_score": 0.7,
                "verdict": "Vorgaben unvollständig",
            },
            "missing_information": ["Abgabedatum", "Bewertungskriterien"],
        }

        md = format_audit_as_md(audit_data)
        self.assertIn("# Vorgaben-Audit", md)
        self.assertIn("Blocker", md)
        self.assertIn("F-001", md)
        self.assertIn("SQLite", md)
        self.assertIn("Fehlende Informationen", md)
        self.assertIn("Abgabedatum", md)
        self.assertIn("40%", md)  # completeness score

    def test_audit_report_docx(self):
        """DOCX report is valid and contains findings."""
        from schulpipeline.audit import format_audit_as_docx

        audit_data = {
            "title": "Test-Audit",
            "documents_analyzed": ["Test.docx"],
            "findings": [
                {
                    "id": "F-001",
                    "category": "gap",
                    "severity": "warning",
                    "title": "Testfinding",
                    "detail": "Testdetail",
                    "sources": ["Test.docx"],
                    "quotes": ["Testzitat"],
                    "recommendation": "Testen",
                }
            ],
            "summary": {
                "total_findings": 1,
                "blockers": 0,
                "warnings": 1,
                "info": 0,
                "completeness_score": 0.8,
                "feasibility_score": 0.9,
                "verdict": "Kleinere Unklarheiten",
            },
            "missing_information": [],
        }

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            path = f.name

        try:
            format_audit_as_docx(audit_data, path)
            self.assertTrue(Path(path).exists())

            from docx import Document

            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
            self.assertIn("Test-Audit", text)
            self.assertIn("Testfinding", text)
        finally:
            Path(path).unlink(missing_ok=True)

    def test_audit_preset(self):
        """Audit preset resolves and selects correct pipeline stages."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.presets import resolve_preset, resolve_quick

        preset = resolve_preset("audit", "programmierung")
        self.assertEqual(preset.output_format, "audit")
        self.assertTrue(preset.output_constraints.get("audit_only"))

        quick = resolve_quick("fiae-audit")
        self.assertEqual(quick.output_format, "audit")

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)
        stages = pipeline._select_stages({"preset": preset})
        stage_names = [s.name for s in stages]
        self.assertEqual(stage_names, ["intake", "classify_docs", "audit"])

    def test_template_mode_includes_audit(self):
        """Template fill flow now includes audit as intermediate step."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.presets import resolve_preset

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)

        preset = resolve_preset("vorlage", "programmierung")
        stages = pipeline._select_stages({"preset": preset})
        stage_names = [s.name for s in stages]
        self.assertEqual(stage_names, ["intake", "classify_docs", "audit", "fill_template"])

    def test_severity_ordering(self):
        """Findings are sorted blockers → warnings → info."""
        from schulpipeline.audit import AuditStage

        stage = AuditStage()
        mixed = [
            {"severity": "info", "title": "A"},
            {"severity": "blocker", "title": "B"},
            {"severity": "warning", "title": "C"},
        ]
        merged = stage._merge_findings(mixed, {"_findings": []})
        severities = [f["severity"] for f in merged]
        self.assertEqual(severities, ["blocker", "warning", "info"])


class TestRequirements(unittest.TestCase):
    """Tests for the three-part requirements documentation."""

    def test_classify_report_stage(self):
        """ClassifyReportStage extracts structured requirements."""
        from schulpipeline.requirements import ClassifyReportStage

        config, router = make_mock_router()
        stage = ClassifyReportStage()
        self.assertEqual(stage.name, "classify_report")

        classify_result = json.loads(DEFAULT_RESPONSES["Dokument-Klassifikator"])
        audit_result = json.loads(DEFAULT_RESPONSES["Anforderungs-Auditor"])
        # Add findings with IDs
        for i, f in enumerate(audit_result.get("findings", [])):
            f["id"] = f"F-{i + 1:03d}"

        context = {
            "classify_docs": classify_result,
            "audit": audit_result,
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"ClassifyReport failed: {result.errors}")

        reqs = result.data.get("requirements", [])
        self.assertGreaterEqual(len(reqs), 1)
        self.assertIn("id", reqs[0])
        self.assertIn("status", reqs[0])

    def test_amendments_stage(self):
        """AmendmentsStage generates decisions for findings."""
        from schulpipeline.requirements import AmendmentsStage

        config, router = make_mock_router()
        stage = AmendmentsStage()

        audit_result = {
            "findings": [
                {
                    "id": "F-001",
                    "severity": "warning",
                    "title": "Test finding",
                    "detail": "Test",
                    "category": "contradiction",
                },
            ],
        }
        classify_report = {
            "requirements": [
                {
                    "id": "REQ-001",
                    "text": "Test req",
                    "source": "test.txt",
                    "category": "functional",
                    "status": "clear",
                    "priority": "must",
                },
            ],
        }

        context = {
            "audit": audit_result,
            "classify_report": classify_report,
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success, f"Amendments failed: {result.errors}")

        amendments = result.data.get("amendments", [])
        self.assertGreaterEqual(len(amendments), 1)
        self.assertIn("decision", amendments[0])
        self.assertIn("reasoning", amendments[0])

    def test_amendments_empty_when_no_issues(self):
        """No amendments needed when all findings are info-level."""
        from schulpipeline.requirements import AmendmentsStage

        config, router = make_mock_router()
        stage = AmendmentsStage()

        context = {
            "audit": {"findings": [{"id": "F-001", "severity": "info", "title": "Minor note"}]},
            "classify_report": {"requirements": []},
        }

        result = asyncio.run(stage.run(context, router, config))
        self.assertTrue(result.success)
        self.assertTrue(result.data.get("all_clear"))

    def test_generate_deviations(self):
        """Deviation log generated from impossibility findings + amendments."""
        from schulpipeline.requirements import generate_deviations

        audit = {
            "findings": [
                {
                    "id": "F-001",
                    "category": "impossibility",
                    "severity": "blocker",
                    "title": "Inhalt passt nicht auf 1 Seite",
                    "detail": "15 Felder, nur 3000 Zeichen verfuegbar. Benoetigt: 4500.",
                    "sources": ["Antrag.docx"],
                },
                {
                    "id": "F-002",
                    "category": "contradiction",  # Not impossibility, should be skipped
                    "severity": "warning",
                    "title": "DB-Widerspruch",
                    "detail": "...",
                    "sources": [],
                },
            ],
        }

        amendments = {
            "amendments": [
                {
                    "id": "AMD-001",
                    "resolves": "F-001",
                    "decision": "Auf 2 Seiten erweitert",
                    "reasoning": "Mathematisch nicht auf 1 Seite machbar",
                },
            ],
        }

        deviations = generate_deviations(audit, amendments)
        self.assertEqual(len(deviations), 1)  # Only impossibilities
        self.assertEqual(deviations[0]["id"], "DEV-001")
        self.assertIn("1 Seite", deviations[0]["constraint"])
        self.assertIn("2 Seiten", deviations[0]["alternative"])
        self.assertEqual(deviations[0]["severity"], "major")

    def test_build_full_report(self):
        """Full three-part report assembles correctly."""
        from schulpipeline.requirements import build_full_report

        classify_report = {
            "requirements": [
                {
                    "id": "REQ-001",
                    "text": "Login-System",
                    "source": "Anf.txt",
                    "category": "functional",
                    "status": "clear",
                    "priority": "must",
                    "quote": "Login",
                },
            ],
            "implicit_requirements": [],
            "requirement_count": {"total": 1, "clear": 1, "ambiguous": 0, "contradicted": 0, "gap": 0},
        }
        audit = {
            "findings": [
                {
                    "id": "F-001",
                    "category": "impossibility",
                    "severity": "blocker",
                    "title": "Seitenconstraint",
                    "detail": "Zu viel fuer 1 Seite",
                    "sources": ["X"],
                },
            ],
            "summary": {
                "total_findings": 1,
                "blockers": 1,
                "warnings": 0,
                "info": 0,
                "completeness_score": 0.4,
                "feasibility_score": 0.7,
                "verdict": "Vorgaben unvollstaendig",
            },
        }
        amendments = {
            "amendments": [
                {
                    "id": "AMD-001",
                    "resolves": "F-001",
                    "decision": "2 Seiten",
                    "reasoning": "Mathe",
                    "source": "auto",
                    "alternatives_considered": [],
                },
            ],
            "unresolvable": [],
        }

        report = build_full_report(classify_report, audit, amendments)
        self.assertIn("part_a", report)
        self.assertIn("part_b", report)
        self.assertIn("part_c", report)
        self.assertEqual(len(report["part_a"]["requirements"]), 1)
        self.assertEqual(len(report["part_b"]["amendments"]), 1)
        self.assertEqual(len(report["part_c"]["deviations"]), 1)

    def test_format_report_md(self):
        """Markdown report contains all three parts."""
        from schulpipeline.requirements import build_full_report, format_report_as_md

        report = build_full_report(
            {
                "requirements": [
                    {
                        "id": "REQ-001",
                        "text": "Test",
                        "source": "X",
                        "category": "functional",
                        "status": "clear",
                        "priority": "must",
                        "quote": "test quote",
                    }
                ],
                "implicit_requirements": [],
                "requirement_count": {"total": 1, "clear": 1},
            },
            {
                "findings": [],
                "summary": {
                    "verdict": "OK",
                    "completeness_score": 1.0,
                    "feasibility_score": 1.0,
                    "total_findings": 0,
                    "blockers": 0,
                    "warnings": 0,
                    "info": 0,
                },
            },
            {"amendments": [], "unresolvable": []},
        )

        md = format_report_as_md(report)
        self.assertIn("Teil A", md)
        self.assertIn("Teil B", md)
        self.assertIn("Teil C", md)
        self.assertIn("REQ-001", md)
        self.assertIn("Keine Abweichungen", md)

    def test_format_report_docx(self):
        """DOCX report is valid."""
        from schulpipeline.requirements import build_full_report, format_report_as_docx

        report = build_full_report(
            {
                "requirements": [
                    {
                        "id": "REQ-001",
                        "text": "Test",
                        "source": "X",
                        "category": "functional",
                        "status": "clear",
                        "priority": "must",
                        "quote": "original",
                    }
                ],
                "implicit_requirements": [],
                "requirement_count": {},
            },
            {"findings": [], "summary": {"verdict": "OK", "completeness_score": 1.0, "feasibility_score": 1.0}},
            {"amendments": [], "unresolvable": []},
        )

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            path = f.name
        try:
            format_report_as_docx(report, path)
            self.assertTrue(Path(path).exists())
            from docx import Document

            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
            self.assertIn("Anforderungsdokumentation", text)
            self.assertIn("Teil A", text)
        finally:
            Path(path).unlink(missing_ok=True)

    def test_requirements_preset_and_stages(self):
        """Requirements report preset selects correct pipeline stages."""
        from schulpipeline.pipeline import Pipeline
        from schulpipeline.presets import resolve_preset, resolve_quick

        preset = resolve_preset("anforderungen", "programmierung")
        self.assertEqual(preset.output_format, "requirements_report")
        self.assertTrue(preset.output_constraints.get("requirements_report"))

        quick = resolve_quick("fiae-anforderungen")
        self.assertEqual(quick.output_format, "requirements_report")

        config, router = make_mock_router()
        pipeline = Pipeline(config, router)
        stages = pipeline._select_stages({"preset": preset})
        stage_names = [s.name for s in stages]
        self.assertEqual(stage_names, ["intake", "classify_docs", "audit", "classify_report", "amendments"])

    def test_cross_reference(self):
        """Requirements are cross-referenced with audit findings."""
        from schulpipeline.requirements import _cross_reference

        requirements = [
            {"id": "REQ-001", "text": "Datenbank SQLite verwenden"},
            {"id": "REQ-002", "text": "Login-System implementieren"},
        ]
        findings = [
            {
                "id": "F-001",
                "title": "Widerspruch Datenbank SQLite vs PostgreSQL",
                "detail": "SQLite in Anforderungen, PostgreSQL muendlich",
            },
        ]

        result = _cross_reference(requirements, findings)
        # REQ-001 should be linked to F-001 (both mention SQLite/Datenbank)
        self.assertGreater(len(result[0].get("related_findings", [])), 0)


class TestFeedback(unittest.TestCase):
    """Tests for feedback collection and research export."""

    def setUp(self):
        self.tmpdir = tempfile.mkdtemp()

    def tearDown(self):
        import shutil

        shutil.rmtree(self.tmpdir, ignore_errors=True)

    def test_save_and_load_record(self):
        """Feedback records persist correctly."""
        from schulpipeline.feedback import FeedbackRecord, FeedbackStore

        store = FeedbackStore(Path(self.tmpdir) / "feedback")
        record = FeedbackRecord(
            run_id="test-001",
            timestamp="2026-01-01T00:00:00Z",
            pipeline_flow="worksheet",
            preset_key="fiae-blatt-wirtschaft",
            subject="wirtschaft",
            output_format="worksheet",
            total_stages=3,
            failed_stages=0,
            total_cost_usd=0.0,
            elapsed_ms=1500,
            backends_used=["groq"],
            quality_rating=4,
            estimated_time_saved_min=25,
        )

        store.save_record(record)
        loaded = store.load_record("test-001")
        self.assertIsNotNone(loaded)
        self.assertEqual(loaded.run_id, "test-001")
        self.assertEqual(loaded.quality_rating, 4)
        self.assertEqual(loaded.estimated_time_saved_min, 25)

    def test_update_record(self):
        """Records can be updated (e.g. adding grade later)."""
        from schulpipeline.feedback import FeedbackRecord, FeedbackStore

        store = FeedbackStore(Path(self.tmpdir) / "feedback")
        record = FeedbackRecord(
            run_id="test-002",
            timestamp="2026-01-01T00:00:00Z",
            pipeline_flow="presentation",
            preset_key="test",
            subject="test",
            output_format="pptx",
            total_stages=5,
            failed_stages=0,
            total_cost_usd=0.0,
            elapsed_ms=2000,
            backends_used=[],
            grade_received="pending",
        )
        store.save_record(record)

        updated = store.update_record("test-002", grade_received="2", quality_rating=4)
        self.assertEqual(updated.grade_received, "2")
        self.assertEqual(updated.quality_rating, 4)

        # Verify persisted
        reloaded = store.load_record("test-002")
        self.assertEqual(reloaded.grade_received, "2")

    def test_aggregate_stats(self):
        """Aggregates computed correctly across records."""
        from schulpipeline.feedback import FeedbackRecord, FeedbackStore

        store = FeedbackStore(Path(self.tmpdir) / "feedback")

        for i in range(5):
            record = FeedbackRecord(
                run_id=f"test-{i:03d}",
                timestamp="2026-01-01T00:00:00Z",
                pipeline_flow="worksheet" if i < 3 else "presentation",
                preset_key="test",
                subject="wirtschaft" if i < 3 else "informatik",
                output_format="worksheet",
                total_stages=3,
                failed_stages=0,
                total_cost_usd=0.0,
                elapsed_ms=1000 + i * 100,
                backends_used=["groq"],
                quality_rating=3 + (i % 3),  # 3, 4, 5, 3, 4
                estimated_time_saved_min=20 + i * 5,
                usable_without_edits=(i % 2 == 0),
            )
            store.save_record(record)

        stats = store.get_aggregates()
        self.assertEqual(stats.total_runs, 5)
        self.assertEqual(stats.total_time_saved_min, 20 + 25 + 30 + 35 + 40)
        self.assertAlmostEqual(stats.pct_free_runs, 1.0)  # All free
        self.assertIn("worksheet", stats.runs_by_flow)
        self.assertEqual(stats.runs_by_flow["worksheet"], 3)
        self.assertGreater(stats.avg_time_saved_min, 0)

    def test_delete_all(self):
        """Delete all wipes everything."""
        from schulpipeline.feedback import FeedbackRecord, FeedbackStore

        store = FeedbackStore(Path(self.tmpdir) / "feedback")
        for i in range(3):
            store.save_record(
                FeedbackRecord(
                    run_id=f"del-{i}",
                    timestamp="2026-01-01T00:00:00Z",
                    pipeline_flow="test",
                    preset_key="",
                    subject="",
                    output_format="",
                    total_stages=1,
                    failed_stages=0,
                    total_cost_usd=0.0,
                    elapsed_ms=100,
                    backends_used=[],
                )
            )

        self.assertEqual(len(store.all_records()), 3)
        deleted = store.delete_all()
        self.assertEqual(deleted, 3)
        self.assertEqual(len(store.all_records()), 0)

    def test_research_export(self):
        """Research export contains only aggregates, no PII."""
        from schulpipeline.feedback import FeedbackRecord, FeedbackStore, export_for_research

        store = FeedbackStore(Path(self.tmpdir) / "feedback")

        # Add records with audit data
        for i in range(3):
            store.save_record(
                FeedbackRecord(
                    run_id=f"exp-{i}",
                    timestamp="2026-01-01T00:00:00Z",
                    pipeline_flow="audit",
                    preset_key="fiae-audit",
                    subject="programmierung",
                    output_format="audit",
                    total_stages=3,
                    failed_stages=0,
                    total_cost_usd=0.0,
                    elapsed_ms=1500,
                    backends_used=["groq"],
                    quality_rating=4,
                    estimated_time_saved_min=30,
                    audit_findings_total=5,
                    audit_blockers=1,
                    audit_warnings=3,
                    audit_completeness=0.4,
                    audit_feasibility=0.7,
                    contradictions_found=2,
                )
            )

        export = export_for_research(store)

        # Verify structure
        self.assertEqual(export["_schema_version"], "1.0")
        self.assertEqual(export["sample_size"], 3)
        self.assertIn("usage", export)
        self.assertIn("quality", export)
        self.assertIn("audit_findings", export)
        self.assertIn("cost", export)

        # Verify audit metrics
        audit = export["audit_findings"]
        self.assertEqual(audit["avg_completeness_pct"], 40.0)
        self.assertEqual(audit["avg_feasibility_pct"], 70.0)
        self.assertEqual(audit["avg_contradictions_per_assignment"], 2.0)
        self.assertEqual(audit["pct_assignments_with_blockers"], 100.0)

        # Verify NO individual data leaked
        export_str = json.dumps(export)
        self.assertNotIn("exp-0", export_str)  # No run IDs
        self.assertNotIn("feedback_text", export_str)  # No free text

    def test_research_export_md(self):
        """Markdown export is readable."""
        from schulpipeline.feedback import FeedbackRecord, FeedbackStore, export_for_research, format_research_export_md

        store = FeedbackStore(Path(self.tmpdir) / "feedback")
        store.save_record(
            FeedbackRecord(
                run_id="md-1",
                timestamp="2026-01-01T00:00:00Z",
                pipeline_flow="worksheet",
                preset_key="test",
                subject="wirtschaft",
                output_format="worksheet",
                total_stages=3,
                failed_stages=0,
                total_cost_usd=0.0,
                elapsed_ms=1000,
                backends_used=["groq"],
                quality_rating=4,
                estimated_time_saved_min=20,
                audit_findings_total=3,
                audit_blockers=1,
                audit_completeness=0.5,
                audit_feasibility=0.8,
                contradictions_found=1,
            )
        )

        export = export_for_research(store)
        md = format_research_export_md(export)
        self.assertIn("Schulpipeline", md)
        self.assertIn("Zeitersparnis", md)
        self.assertIn("Aufgabenqualität", md)
        self.assertIn("Anonymisierte", md)

    def test_build_feedback_from_result(self):
        """Feedback auto-built from pipeline result."""
        from schulpipeline.feedback import build_feedback_from_result
        from schulpipeline.stages.base import StageResult

        # Mock a pipeline result
        class MockResult:
            success = True
            results = [
                StageResult(stage="intake", success=True, data={}),
                StageResult(stage="decompose", success=True, data={}),
                StageResult(stage="solve", success=True, data={}),
            ]
            total_cost_usd = 0.0
            elapsed_ms = 1200

        record = build_feedback_from_result("session-abc", MockResult())
        self.assertEqual(record.pipeline_flow, "worksheet")
        self.assertEqual(record.total_stages, 3)
        self.assertEqual(record.failed_stages, 0)
        self.assertEqual(record.elapsed_ms, 1200)

    def test_instance_id_stable(self):
        """Same store produces same instance ID."""
        from schulpipeline.feedback import FeedbackStore, _generate_instance_id

        store = FeedbackStore(Path(self.tmpdir) / "feedback")
        id1 = _generate_instance_id(store)
        id2 = _generate_instance_id(store)
        self.assertEqual(id1, id2)
        self.assertEqual(len(id1), 12)


if __name__ == "__main__":
    print(f"Python {sys.version}")
    print("Running schulpipeline tests...\n")
    unittest.main(verbosity=2)
