"""Tests for artifact builders — PPTX, DOCX, MD file generation."""

import pytest

from schulpipeline.artifacts.converter import synthesis_to_presentation
from schulpipeline.artifacts.docx_builder import build_docx
from schulpipeline.artifacts.md_builder import build_md
from schulpipeline.artifacts.pptx_builder import build_pptx
from schulpipeline.artifacts.slide_registry import (
    _REGISTRY,
    _get_spec,
    classify_presentation,
    classify_section,
)
from schulpipeline.presets import resolve_preset

# ============================================================
# PPTX Builder
# ============================================================


def test_pptx_creates_file(synthesis_data, tmp_path):
    out = tmp_path / "test.pptx"
    build_pptx(synthesis_data, out)
    assert out.exists()
    assert out.stat().st_size > 0


def test_pptx_has_correct_slide_count(synthesis_data, tmp_path):
    out = tmp_path / "slides.pptx"
    build_pptx(synthesis_data, out)

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == len(synthesis_data["sections"])


def test_pptx_title_slide_has_heading(synthesis_data, tmp_path):
    out = tmp_path / "title.pptx"
    build_pptx(synthesis_data, out)

    from pptx import Presentation

    prs = Presentation(str(out))
    first_slide = prs.slides[0]
    texts = [shape.text for shape in first_slide.shapes if shape.has_text_frame]
    assert any(synthesis_data["sections"][0]["heading"] in t for t in texts)


def test_pptx_minimal_input(tmp_path):
    """Builder handles minimal synthesis data."""
    data = {
        "title": "Minimal",
        "sections": [
            {
                "section_id": "s1",
                "heading": "Only Slide",
                "content": "Content here",
                "bullet_points": [],
                "speaker_notes": None,
            }
        ],
        "sources": [],
    }
    out = tmp_path / "minimal.pptx"
    build_pptx(data, out)
    assert out.exists()

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_pptx_empty_section_skipped(tmp_path):
    """Sections with no heading, content, or bullets produce no slide."""
    data = {
        "title": "Test",
        "sections": [
            {
                "section_id": "s1",
                "heading": "Real Slide",
                "content": "Content",
                "bullet_points": ["A"],
                "speaker_notes": None,
            },
            {"section_id": "s2", "heading": "", "content": "", "bullet_points": [], "speaker_notes": None},
        ],
        "sources": [],
    }
    out = tmp_path / "skip_empty.pptx"
    build_pptx(data, out)

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_pptx_sources_detection_variants(tmp_path):
    """Various sources heading names are all detected as sources slides."""
    for heading in ["Quellen", "Quellenangaben", "Referenzen", "Literaturverzeichnis", "Sources"]:
        data = {
            "title": "Test",
            "sections": [
                {
                    "section_id": "s1",
                    "heading": "Content",
                    "content": "Text here",
                    "bullet_points": ["A"],
                    "speaker_notes": None,
                },
                {
                    "section_id": "s2",
                    "heading": heading,
                    "content": "",
                    "bullet_points": ["Source 1", "Source 2"],
                    "speaker_notes": None,
                },
            ],
            "sources": [],
        }
        out = tmp_path / f"sources_{heading}.pptx"
        build_pptx(data, out)

        from pptx import Presentation

        prs = Presentation(str(out))
        # Should have 2 slides (content + sources), sources detected properly
        assert len(prs.slides) == 2


# ============================================================
# DOCX Builder — Sources Deduplication
# ============================================================


def test_docx_sources_dedup(tmp_path):
    """Sources from inline sections and top-level are deduplicated."""
    data = {
        "title": "Test",
        "sections": [
            {"heading": "Content", "content": "Some text.", "bullet_points": []},
            {"heading": "Quellen", "content": "", "bullet_points": ["Source A", "Source B"]},
        ],
        "sources": ["Source A", "Source C"],
    }
    out = tmp_path / "dedup.docx"
    build_docx(data, out)

    from docx import Document

    doc = Document(str(out))
    all_text = [p.text for p in doc.paragraphs]
    # "Source A" should appear only once (deduplicated)
    assert all_text.count("Source A") == 1
    assert "Source B" in all_text
    assert "Source C" in all_text


def test_docx_no_sources_when_empty(tmp_path):
    """No Quellen heading when there are no sources at all."""
    data = {
        "title": "Bare",
        "sections": [{"heading": "Chapter", "content": "Text.", "bullet_points": []}],
        "sources": [],
    }
    out = tmp_path / "nosources.docx"
    build_docx(data, out)

    from docx import Document

    doc = Document(str(out))
    headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
    assert "Quellen" not in headings


# ============================================================
# DOCX Builder
# ============================================================


def test_docx_creates_file(synthesis_data, tmp_path):
    out = tmp_path / "test.docx"
    build_docx(synthesis_data, out)
    assert out.exists()
    assert out.stat().st_size > 0


def test_docx_contains_title(synthesis_data, tmp_path):
    out = tmp_path / "title.docx"
    build_docx(synthesis_data, out)

    from docx import Document

    doc = Document(str(out))
    paragraphs = [p.text for p in doc.paragraphs]
    assert synthesis_data["title"] in paragraphs


def test_docx_contains_sources(synthesis_data, tmp_path):
    out = tmp_path / "sources.docx"
    build_docx(synthesis_data, out)

    from docx import Document

    doc = Document(str(out))
    all_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Quellen" in all_text


def test_docx_minimal_input(tmp_path):
    data = {
        "title": "Minimal Doc",
        "sections": [{"heading": "Chapter 1", "content": "Some text here.", "bullet_points": []}],
        "sources": [],
    }
    out = tmp_path / "minimal.docx"
    build_docx(data, out)
    assert out.exists()


# ============================================================
# Markdown Builder
# ============================================================


def test_md_creates_file(synthesis_data, tmp_path):
    out = tmp_path / "test.md"
    build_md(synthesis_data, out)
    assert out.exists()
    assert out.stat().st_size > 0


def test_md_has_title_heading(synthesis_data, tmp_path):
    out = tmp_path / "heading.md"
    build_md(synthesis_data, out)

    content = out.read_text(encoding="utf-8")
    assert content.startswith(f"# {synthesis_data['title']}")


def test_md_has_section_headings(synthesis_data, tmp_path):
    out = tmp_path / "sections.md"
    build_md(synthesis_data, out)

    content = out.read_text(encoding="utf-8")
    for section in synthesis_data["sections"]:
        assert f"## {section['heading']}" in content


def test_md_has_bullet_points(synthesis_data, tmp_path):
    out = tmp_path / "bullets.md"
    build_md(synthesis_data, out)

    content = out.read_text(encoding="utf-8")
    # Section 2 has bullet points
    assert "- Schutz von Informationen" in content


def test_md_has_sources(synthesis_data, tmp_path):
    out = tmp_path / "sources.md"
    build_md(synthesis_data, out)

    content = out.read_text(encoding="utf-8")
    assert "## Quellen" in content
    assert "- BSI" in content


def test_md_minimal_input(tmp_path):
    data = {"title": "Minimal", "sections": [], "sources": []}
    out = tmp_path / "minimal.md"
    build_md(data, out)
    assert out.exists()
    content = out.read_text(encoding="utf-8")
    assert "# Minimal" in content


# ============================================================
# PPTX Builder — Backward Compatibility
# ============================================================


def test_pptx_backward_compat(synthesis_data, tmp_path):
    """build_pptx() without preset still works (defaults)."""
    out = tmp_path / "compat.pptx"
    build_pptx(synthesis_data, out)
    assert out.exists()
    assert out.stat().st_size > 0


# ============================================================
# DOCX Builder — Style & Visual Tests
# ============================================================


def test_docx_backward_compat(synthesis_data, tmp_path):
    """build_docx() without style args still works (defaults)."""
    out = tmp_path / "compat.docx"
    build_docx(synthesis_data, out)
    assert out.exists()
    assert out.stat().st_size > 0


def test_docx_with_style_font(synthesis_data, modern_style, visual_config_disabled, tmp_path):
    """Modern style sets font to Arial."""
    out = tmp_path / "modern.docx"
    build_docx(synthesis_data, out, modern_style.visual, visual_config_disabled)

    from docx import Document

    doc = Document(str(out))
    normal_font = doc.styles["Normal"].font
    assert normal_font.name == "Arial"


def test_docx_visual_placeholder(synthesis_data_with_visuals, clean_style, visual_config_enabled, tmp_path):
    """Visual placeholder renders as an italic paragraph in the document."""
    out = tmp_path / "visuals.docx"
    build_docx(synthesis_data_with_visuals, out, clean_style.visual, visual_config_enabled)

    from docx import Document

    doc = Document(str(out))
    all_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Angriffsarten" in all_text  # The intent text should appear


def test_docx_visuals_disabled(synthesis_data_with_visuals, clean_style, visual_config_disabled, tmp_path):
    """No placeholder when visuals are disabled."""
    out = tmp_path / "no_visuals.docx"
    build_docx(synthesis_data_with_visuals, out, clean_style.visual, visual_config_disabled)

    from docx import Document

    doc = Document(str(out))
    # Check that the placeholder intent text does NOT appear
    italic_paragraphs = [p for p in doc.paragraphs if p.runs and any(r.italic for r in p.runs)]
    # No italic visual placeholder paragraphs
    assert len(italic_paragraphs) == 0


# ============================================================
# Slide Registry — classify_section
# ============================================================


class TestClassifySection:
    def _section(self, heading="", content="", bullets=None, visuals=None):
        return {
            "heading": heading,
            "content": content,
            "bullet_points": bullets or [],
            "visuals": visuals or [],
        }

    # --- Position constraints ---

    def test_first_section_always_title(self):
        """Index 0 always returns the title spec regardless of content."""
        section = self._section(heading="Intro", bullets=["a", "b", "c"])
        spec = classify_section(section, index=0, total=5)
        assert spec.key == "title"
        assert spec.render_fn == "title"

    def test_section_break_never_last(self):
        """section_break spec has never_last=True — should not win at last index."""
        section = self._section(heading="Abschnitt 2")
        spec = classify_section(
            section,
            index=4,
            total=5,
            structure_hint="section_break",
            style="bullet-heavy",
        )
        assert spec.key != "section_break"

    def test_sources_keyword_overrides_score(self):
        """A section whose heading contains a sources keyword wins the sources spec."""
        for heading in ["Quellen", "quellenangaben", "Literaturverzeichnis", "Sources"]:
            section = self._section(heading=heading, bullets=["src1", "src2"])
            spec = classify_section(section, index=3, total=5)
            assert spec.key == "sources", f"Expected sources for heading={heading!r}"

    # --- Structure hint bonus ---

    def test_structure_hint_fazit_prefers_closing(self):
        """Structure hint 'fazit' should produce the closing spec."""
        section = self._section(
            heading="Fazit",
            content="Zusammenfassend lässt sich sagen...",
        )
        spec = classify_section(
            section,
            index=3,
            total=5,
            structure_hint="fazit",
            style="prose",
        )
        assert spec.key == "closing"

    def test_structure_hint_inhalt_prefers_content(self):
        """Structure hint 'inhalt' with bullets should resolve to content."""
        section = self._section(
            heading="Firewall",
            bullets=["Punkt 1", "Punkt 2", "Punkt 3"],
        )
        spec = classify_section(
            section,
            index=2,
            total=5,
            structure_hint="inhalt",
            style="bullet-heavy",
        )
        assert spec.key == "content"

    def test_unknown_structure_hint_still_scores(self):
        """An unrecognised structure hint does not crash — scoring proceeds normally."""
        section = self._section(heading="Irgendwas", bullets=["a", "b"])
        spec = classify_section(
            section,
            index=1,
            total=5,
            structure_hint="unbekannt",
            style="bullet-heavy",
        )
        assert spec.key in {s.key for s in _REGISTRY}

    # --- Style multipliers ---

    def test_prose_style_prefers_body_over_bullets(self):
        """In prose style, a section with only body text should not be section_break."""
        body_section = self._section(
            heading="Einleitung",
            content=" ".join(["word"] * 60),
        )
        spec = classify_section(
            body_section,
            index=1,
            total=5,
            structure_hint="einleitung",
            style="prose",
        )
        assert spec.key in ("intro", "content", "closing")
        assert spec.key != "section_break"

    # --- Fallback behavior ---

    def test_no_preset_backward_compat(self):
        """No structure hint, no style → classify_section still returns a valid spec."""
        section = self._section(heading="IT-Sicherheit", bullets=["CIA", "BSI"])
        spec = classify_section(section, index=1, total=5)
        assert spec is not None
        assert spec.render_fn in ("title", "content", "sources", "section_break")


class TestClassifyPresentation:
    def _sections(self, count=5):
        return [
            {
                "heading": f"Section {i}",
                "content": "text",
                "bullet_points": ["a", "b"],
                "visuals": [],
                "speaker_notes": None,
            }
            for i in range(count)
        ]

    def test_returns_one_spec_per_section(self):
        sections = self._sections(5)
        result = classify_presentation(sections)
        assert len(result) == len(sections)

    def test_first_is_always_title(self):
        sections = self._sections(5)
        result = classify_presentation(sections)
        assert result[0].key == "title"

    def test_with_preset_uses_structure(self):
        preset = resolve_preset("praesentation", "it_sicherheit")
        sections = self._sections(len(preset.structure))
        result = classify_presentation(sections, preset=preset)
        # Last structure hint is "quellen" — should classify as sources
        assert result[-1].key == "sources"

    def test_no_preset_returns_valid_specs(self):
        sections = self._sections(3)
        result = classify_presentation(sections, preset=None)
        for spec in result:
            assert spec.key in {s.key for s in _REGISTRY}

    def test_empty_sections_returns_empty(self):
        result = classify_presentation([])
        assert result == []


class TestRegistryIntegrity:
    """Sanity checks on the registry definition itself."""

    def test_all_render_fns_are_known(self):
        known = {"title", "content", "sources", "section_break"}
        for spec in _REGISTRY:
            assert spec.render_fn in known, f"{spec.key} has unknown render_fn={spec.render_fn}"

    def test_exactly_one_force_at_first(self):
        forced = [s for s in _REGISTRY if s.force_at_first]
        assert len(forced) == 1
        assert forced[0].key == "title"

    def test_structure_aliases_are_frozensets(self):
        for spec in _REGISTRY:
            assert isinstance(spec.structure_aliases, frozenset), f"{spec.key}.structure_aliases is not a frozenset"

    def test_get_spec_raises_on_unknown_key(self):
        with pytest.raises(KeyError):
            _get_spec("nonexistent_key")


# ============================================================
# Converter — synthesis_to_presentation
# ============================================================


class TestSynthesisToPresentation:
    def test_basic_conversion(self, synthesis_data):
        """Synthesis dict converts to a slideforge Presentation with correct slide count."""
        pres = synthesis_to_presentation(synthesis_data)
        non_empty = [
            s for s in synthesis_data["sections"] if s.get("heading") or s.get("content") or s.get("bullet_points")
        ]
        assert len(pres.slides) == len(non_empty)
        assert pres.name == synthesis_data["title"]

    def test_title_slide_uses_sp_title_layout(self, synthesis_data):
        """First slide gets SP_Title layout."""
        pres = synthesis_to_presentation(synthesis_data)
        assert pres.slides[0].layout == "SP_Title"

    def test_sources_section_gets_sp_sources_layout(self):
        """A section with a sources heading maps to SP_Sources."""
        data = {
            "title": "Test",
            "sections": [
                {"heading": "Intro", "content": "Welcome", "bullet_points": []},
                {"heading": "Quellen", "content": "", "bullet_points": ["Src 1"]},
            ],
            "sources": ["Src 1", "Src 2"],
        }
        pres = synthesis_to_presentation(data)
        sources_slide = pres.slides[-1]
        assert sources_slide.layout == "SP_Sources"
        # Sources body uses top-level sources list
        assert "Src 1" in sources_slide.body
        assert "Src 2" in sources_slide.body

    def test_empty_sections_skipped(self):
        """Sections with no heading, content, or bullets produce no slide."""
        data = {
            "title": "Test",
            "sections": [
                {"heading": "Real", "content": "Content", "bullet_points": []},
                {"heading": "", "content": "", "bullet_points": []},
            ],
            "sources": [],
        }
        pres = synthesis_to_presentation(data)
        assert len(pres.slides) == 1

    def test_bullets_as_newline_body(self):
        """Bullet points are joined with newlines in the body."""
        data = {
            "title": "Test",
            "sections": [
                {"heading": "Title", "content": "Sub", "bullet_points": []},
                {"heading": "Points", "content": "", "bullet_points": ["Alpha", "Beta", "Gamma"]},
            ],
            "sources": [],
        }
        pres = synthesis_to_presentation(data)
        content_slide = pres.slides[1]
        assert content_slide.body == "Alpha\nBeta\nGamma"

    def test_speaker_notes_preserved(self):
        """Speaker notes from synthesis are passed through."""
        data = {
            "title": "Test",
            "sections": [
                {"heading": "Slide", "content": "Body", "bullet_points": [], "speaker_notes": "Say this aloud"},
            ],
            "sources": [],
        }
        pres = synthesis_to_presentation(data)
        assert pres.slides[0].notes == "Say this aloud"

    def test_custom_name_overrides_title(self):
        """The name parameter overrides the synthesis title."""
        data = {"title": "Original", "sections": [], "sources": []}
        pres = synthesis_to_presentation(data, name="Custom Name")
        assert pres.name == "Custom Name"
